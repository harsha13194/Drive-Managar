# app.py - Live Google Drive Integration for Drive Manager
import sys
try:
    import requests
except ImportError as e:
    print("Startup Validation Error: requests module is missing. Please run 'pip install requests'.", file=sys.stderr)
    sys.exit(1)

import logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

from flask import Flask, render_template, request, jsonify, send_file, redirect, url_for, session, g, has_request_context
import os
import json
from datetime import datetime
import io
import re
import threading
import difflib
import math
import traceback
import queue
import sqlite3
import embeddings_db
import document_tools

# NOTE: No hardcoded fallback API key is used. Configure NVIDIA_API_KEY (and/or
# GEMINI_API_KEY) as environment variables, or let users supply their own key
# via the session ("set key ..." chat command in the AI Assistant).
DEFAULT_NVIDIA_API_KEY = os.environ.get('NVIDIA_API_KEY')

# Thread-safe lock for reading/writing JSON indexes
index_lock = threading.Lock()

# Global in-memory indexing status storage
INDEXING_STATUS = {}

# Thread-safe queue for AI file processing
AI_PROCESSING_QUEUE = queue.Queue()

class FileContext:
    def __init__(self, selected_file_id=None, selected_file_name=None, selected_file_type=None):
        self.selected_file_id = selected_file_id
        self.selected_file_name = selected_file_name
        self.selected_file_type = selected_file_type

    def to_dict(self):
        return {
            "selected_file_id": self.selected_file_id,
            "selected_file_name": self.selected_file_name,
            "selected_file_type": self.selected_file_type
        }

# Import Google Libraries
import google.oauth2.credentials
import google_auth_oauthlib.flow
from googleapiclient.discovery import build
from googleapiclient.http import MediaIoBaseUpload, MediaIoBaseDownload

# Allow insecure transport for localhost OAuth (HTTP instead of HTTPS)
os.environ['OAUTHLIB_INSECURE_TRANSPORT'] = '1'
os.environ['OAUTHLIB_RELAX_TOKEN_SCOPE'] = '1'

app = Flask(__name__)
# Flask session secret: prefer an explicit environment variable so sessions
# survive process restarts; otherwise fall back to a securely generated key.
app.secret_key = os.environ.get('FLASK_SECRET_KEY') or os.urandom(32).hex()
app.config['MAX_CONTENT_LENGTH'] = 100 * 1024 * 1024  # 100MB max upload

# OAuth Scopes
SCOPES = [
    'https://www.googleapis.com/auth/drive',
    'openid',
    'https://www.googleapis.com/auth/userinfo.profile',
    'https://www.googleapis.com/auth/userinfo.email'
]


def create_oauth_flow(redirect_uri=None):
    """Create a google_auth_oauthlib.flow.Flow using environment variables.

    Priority: use `GOOGLE_CLIENT_ID` and `GOOGLE_CLIENT_SECRET` environment
    variables. If those are missing, attempt to load `client_secrets.json` if
    present; otherwise raise a clear error.
    """
    client_id = os.environ.get('GOOGLE_CLIENT_ID')
    client_secret = os.environ.get('GOOGLE_CLIENT_SECRET')
    if client_id and client_secret:
        client_config = {
            "web": {
                "client_id": client_id,
                "client_secret": client_secret,
                "auth_uri": "https://accounts.google.com/o/oauth2/auth",
                "token_uri": "https://oauth2.googleapis.com/token",
                "auth_provider_x509_cert_url": "https://www.googleapis.com/oauth2/v1/certs",
                "redirect_uris": [redirect_uri]
            }
        }
        return google_auth_oauthlib.flow.Flow.from_client_config(client_config, scopes=SCOPES, redirect_uri=redirect_uri)

    # Fallback: try to load a client_secrets.json file if available (not recommended for public repos)
    if os.path.exists('client_secrets.json'):
        return google_auth_oauthlib.flow.Flow.from_client_secrets_file('client_secrets.json', scopes=SCOPES, redirect_uri=redirect_uri)

    raise RuntimeError('Google OAuth client configuration missing. Set GOOGLE_CLIENT_ID and GOOGLE_CLIENT_SECRET environment variables.')

def credentials_to_dict(credentials):
    """Serialize credentials object to dictionary for session storage (optimized for size)"""
    return {
        'token': credentials.token,
        'refresh_token': credentials.refresh_token,
        'scopes': credentials.scopes
    }

def login_user(credentials):
    """Establish a new user session by fetching their email, clearing old session state,
    and storing the new credentials and email.
    """
    service = build('drive', 'v3', credentials=credentials, static_discovery=True)
    user_email = get_user_email(service)
    
    # Completely clear the session before storing the new user's state to prevent leaks
    session.clear()
    session['credentials'] = credentials_to_dict(credentials)
    session['user_email'] = user_email
    session.modified = True
    return user_email

def get_credentials():
    """Retrieve credentials from Flask session"""
    if 'credentials' not in session:
        return None
    user_email = session.get('user_email')
    if not user_email or user_email == 'default_user':
        return None
        
    creds_dict = session['credentials'].copy()
    creds_dict['token_uri'] = "https://oauth2.googleapis.com/token"
    creds_dict['client_id'] = os.environ.get('GOOGLE_CLIENT_ID')
    creds_dict['client_secret'] = os.environ.get('GOOGLE_CLIENT_SECRET')
    
    if not creds_dict['client_id'] or not creds_dict['client_secret']:
        # Fallback to client_secrets.json if env vars are missing
        if os.path.exists('client_secrets.json'):
            try:
                with open('client_secrets.json', 'r') as f:
                    data = json.load(f)
                    web_config = data.get('web', {})
                    creds_dict['client_id'] = web_config.get('client_id')
                    creds_dict['client_secret'] = web_config.get('client_secret')
            except Exception:
                pass
                
    return google.oauth2.credentials.Credentials(**creds_dict)

def get_drive_service():
    """Build and return an authorized Drive Service client"""
    if has_request_context() and hasattr(g, 'drive_service'):
        return g.drive_service

    creds = get_credentials()
    if not creds:
        return None
    
    # Auto-refresh expired tokens
    if creds.expired and creds.refresh_token:
        try:
            from google.auth.transport.requests import Request
            creds.refresh(Request())
            session['credentials'] = credentials_to_dict(creds)
        except Exception as e:
            logger.error(f"Error refreshing Google OAuth token: {e}")
            return None
            
    service = build('drive', 'v3', credentials=creds, static_discovery=True)
    if has_request_context():
        g.drive_service = service
    return service

def map_drive_file(f):
    """Helper to transform Google Drive API file details into standard App format"""
    mime = f.get('mimeType', '')
    is_folder = mime == 'application/vnd.google-apps.folder'
    
    # Parse owner display name
    owner = 'me'
    if f.get('owners'):
        owner = f.get('owners')[0].get('displayName', 'me')

    # Date formatting
    dt_str = f.get('modifiedTime', '')
    formatted_date = ''
    if dt_str:
        try:
            # Google date: 2025-06-15T09:30:00.000Z
            clean_dt = dt_str.replace('Z', '')
            if '.' in clean_dt:
                clean_dt = clean_dt.split('.')[0]
            dt = datetime.strptime(clean_dt, "%Y-%m-%dT%H:%M:%S")
            formatted_date = dt.strftime("%d %b %Y")
        except Exception:
            formatted_date = dt_str[:10]

    # Format file size
    size_bytes = f.get('size')
    size_int = 0
    size_str = '--'
    if size_bytes:
        size_int = int(size_bytes)
        if size_int < 1024:
            size_str = f"{size_int} B"
        elif size_int < 1024 * 1024:
            size_str = f"{size_int / 1024:.1f} KB"
        else:
            size_str = f"{size_int / (1024 * 1024):.1f} MB"

    # Map file mimeTypes to FontAwesome classes
    icon = 'fa-file'
    type_str = 'file'
    if is_folder:
        icon = 'fa-folder'
        type_str = 'folder'
    elif 'spreadsheet' in mime:
        icon = 'fa-table'
        type_str = 'spreadsheet'
    elif 'presentation' in mime:
        icon = 'fa-chalkboard'
        type_str = 'presentation'
    elif 'document' in mime or 'text/' in mime:
        icon = 'fa-file-alt'
        type_str = 'document'
    elif 'image/' in mime:
        icon = 'fa-image'
        type_str = 'image'
    elif 'pdf' in mime:
        icon = 'fa-file-pdf'
        type_str = 'pdf'
    elif 'video/' in mime:
        icon = 'fa-video'
        type_str = 'video'
    elif 'audio/' in mime:
        icon = 'fa-music'
        type_str = 'audio'
    elif 'zip' in mime or 'archive' in mime:
        icon = 'fa-file-archive'
        type_str = 'archive'

    # Lookup indexing status
    embedding_status = None
    try:
        user_email = session.get('user_email')
        if user_email:
            if not hasattr(g, 'user_index'):
                sanitized_email = sanitize_filename(user_email)
                g.user_index = load_json_file(f"drive_index_{sanitized_email}.json", {})
            file_meta = g.user_index.get(f.get('id'), {})
            embedding_status = file_meta.get("embedding_status")
    except Exception as ex:
        print(f"Error looking up file status: {ex}")

    return {
        "id": f.get('id'),
        "name": f.get('name'),
        "type": type_str,
        "mime_type": mime,
        "icon": icon,
        "size": size_str,
        "size_bytes": size_int,
        "date": formatted_date,
        "modified_time_raw": f.get('modifiedTime', ''),
        "created_time_raw": f.get('createdTime') or f.get('modifiedTime', ''),
        "owner": owner,
        "starred": f.get('starred', False),
        "is_folder": is_folder,
        "webViewLink": f.get('webViewLink'),
        "webContentLink": f.get('webContentLink'),
        "embedding_status": embedding_status
    }

# ========= ROUTING =========

@app.route('/')
def index():
    """Render page"""
    # OAuth callback handling moved to `/oauth2callback` route.

    creds = get_credentials()
    authenticated = False
    if creds:
        authenticated = True
        try:
            service = get_drive_service()
            if service:
                user_email = session.get('user_email')
                if not user_email:
                    user_email = get_user_email(service)
                    session['user_email'] = user_email
                
                # Rebuild and restore user-specific state (skip redundant json parsing/db setups)
                restore_user_state(user_email)
                
                status = INDEXING_STATUS.get(user_email, {}).get("status")
                if not status or status == "failed":
                    creds_dict = credentials_to_dict(creds)
                    nvidia_key = session.get('nvidia_api_key') or os.environ.get('NVIDIA_API_KEY') or DEFAULT_NVIDIA_API_KEY
                    keys_dict = {
                        'nvidia_api_key': nvidia_key
                    }
                    threading.Thread(
                        target=run_background_indexing,
                        args=(creds_dict, user_email, keys_dict),
                        daemon=True
                    ).start()
        except Exception as e:
            print(f"Error checking background indexing status: {e}")
            session.clear()
            authenticated = False
        
    return render_template('drive.html', authenticated=authenticated)

@app.route('/api/indexing-status')
def get_indexing_status():
    user_email = session.get('user_email')
    if not user_email:
        service = get_drive_service()
        if service:
            user_email = get_user_email(service)
            session['user_email'] = user_email
        else:
            return jsonify({"status": "unauthorized"}), 401
            
    status_info = INDEXING_STATUS.get(user_email, {"status": "not_started"})
    return jsonify(status_info)

@app.route('/login')
def login():
    """Start Google Drive OAuth flow"""
    session.clear()  # Clear existing session before creating a new login flow
    flow = create_oauth_flow(redirect_uri=request.url_root.rstrip("/") + "/oauth2callback")
    authorization_url, state = flow.authorization_url(
        access_type='offline',
        include_granted_scopes='true',
        prompt='consent'
    )
    session['state'] = state
    # Preserve PKCE code_verifier for token exchange in the callback
    code_verifier = getattr(flow, 'code_verifier', None)
    if code_verifier:
        session['code_verifier'] = code_verifier
    return redirect(authorization_url)


@app.route('/oauth2callback')
def oauth2callback():
    """OAuth2 callback handler."""
    try:
        flow = create_oauth_flow(redirect_uri=request.url_root.rstrip("/") + "/oauth2callback")
        # Restore PKCE code_verifier saved during /login
        code_verifier = session.get('code_verifier')
        if code_verifier:
            flow.code_verifier = code_verifier
        flow.fetch_token(authorization_response=request.url)
        
        # Log in the user securely
        login_user(flow.credentials)
        return redirect(url_for('index'))
    except Exception as e:
        print(f"Callback authentication failure: {e}")
        return f"Authentication Error: {e}", 400

@app.route('/logout')
def logout():
    """Log user out and clear credentials session"""
    session.clear()
    if hasattr(g, 'user_index'):
        delattr(g, 'user_index')
    response = redirect(url_for('index'))
    response.set_cookie('session', '', expires=0)
    return response

# ========= API ENDPOINTS (LIVE GOOGLE DRIVE) =========

@app.route('/api/files')
def get_files():
    """List items inside a parent folder"""
    parent_id = request.args.get('folder', 'root')
    service = get_drive_service()
    if not service:
        return jsonify({"error": "Unauthorized"}), 401
    
    try:
        query = f"'{parent_id}' in parents and trashed = false"
        results = service.files().list(
            q=query,
            fields="files(id, name, mimeType, size, modifiedTime, createdTime, starred, owners, webViewLink, webContentLink)",
            pageSize=100
        ).execute()
        drive_files = results.get('files', [])
        return jsonify([map_drive_file(f) for f in drive_files])
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route('/api/starred')
def get_starred():
    """List starred items"""
    service = get_drive_service()
    if not service:
        return jsonify({"error": "Unauthorized"}), 401
        
    try:
        query = "starred = true and trashed = false"
        results = service.files().list(
            q=query,
            fields="files(id, name, mimeType, size, modifiedTime, createdTime, starred, owners, webViewLink, webContentLink)",
            pageSize=100
        ).execute()
        drive_files = results.get('files', [])
        return jsonify([map_drive_file(f) for f in drive_files])
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route('/api/recent')
def get_recent():
    """List recently modified documents"""
    service = get_drive_service()
    if not service:
        return jsonify({"error": "Unauthorized"}), 401
        
    try:
        query = "trashed = false and mimeType != 'application/vnd.google-apps.folder'"
        results = service.files().list(
            q=query,
            orderBy="modifiedTime desc",
            fields="files(id, name, mimeType, size, modifiedTime, createdTime, starred, owners, webViewLink, webContentLink)",
            pageSize=15
        ).execute()
        drive_files = results.get('files', [])
        return jsonify([map_drive_file(f) for f in drive_files])
    except Exception as e:
        return jsonify({"error": str(e)}), 500



@app.route('/api/trash')
def get_trash():
    """List trashed files"""
    service = get_drive_service()
    if not service:
        return jsonify({"error": "Unauthorized"}), 401
        
    try:
        query = "trashed = true"
        results = service.files().list(
            q=query,
            fields="files(id, name, mimeType, size, modifiedTime, createdTime, starred, owners, webViewLink, webContentLink)",
            pageSize=100
        ).execute()
        drive_files = results.get('files', [])
        return jsonify([map_drive_file(f) for f in drive_files])
    except Exception as e:
        return jsonify({"error": str(e)}), 500

def get_session_email_or_fetch(service):
    email = session.get('user_email')
    if not email:
        email = get_user_email(service)
        session['user_email'] = email
    return email


def verify_file_ownership(service, user_email, file_id, index_data):
    """
    Verify if the file_id belongs to the user by checking index_data or querying Google Drive API directly.
    If the file exists in Google Drive, dynamically index it to prevent false-positive 403 errors.
    Returns (True, file_metadata) if ownership/access is verified, else (False, None).
    """
    if file_id in index_data:
        return True, index_data[file_id]
        
    try:
        drive_meta = service.files().get(
            fileId=file_id, 
            fields='id, name, mimeType, size, modifiedTime, createdTime, starred, parents, owners, webViewLink, webContentLink'
        ).execute()
        if drive_meta:
            parent_id = drive_meta.get('parents', ['root'])[0] if drive_meta.get('parents') else 'root'
            owner = drive_meta.get('owners')[0].get('displayName', 'me') if drive_meta.get('owners') else 'me'
            file_meta = {
                "id": file_id,
                "name": drive_meta.get("name", "Document"),
                "mime_type": drive_meta.get("mimeType", ""),
                "path": "/" + drive_meta.get("name", "Document"),
                "parent_id": parent_id,
                "size": drive_meta.get("size"),
                "modified_time": drive_meta.get("modifiedTime"),
                "starred": drive_meta.get("starred", False),
                "owner": owner,
                "is_folder": drive_meta.get("mimeType") == 'application/vnd.google-apps.folder',
                "webViewLink": drive_meta.get("webViewLink"),
                "webContentLink": drive_meta.get("webContentLink")
            }
            update_index_file_metadata_only(user_email, file_id, file_meta)
            index_data[file_id] = file_meta
            return True, file_meta
    except Exception as e:
        logger.warning(f"File ownership verification fallback failed for {file_id}: {e}")
        
    return False, None

@app.route('/api/save-doc', methods=['POST'])
def save_doc():
    """Overwrite text file contents on Google Drive"""
    data = request.get_json(silent=True) or {}
    file_id = data.get('id')
    content = data.get('content', '')

    if not file_id:
        return jsonify({"error": "File id is required"}), 400

    service = get_drive_service()
    if not service:
        return jsonify({"error": "Unauthorized"}), 401
        
    try:
        user_email = get_session_email_or_fetch(service)
        if not user_email:
            return jsonify({"error": "Unauthorized"}), 401
        sanitized_email = sanitize_filename(user_email)
        index_data = load_json_file(f"drive_index_{sanitized_email}.json", {})
        has_access, file_meta = verify_file_ownership(service, user_email, file_id, index_data)
        if not has_access:
            return jsonify({"error": "Forbidden - File ownership verification failed"}), 403
            
        media = MediaIoBaseUpload(
            io.BytesIO(content.encode('utf-8')),
            mimetype='text/plain',
            resumable=True
        )
        service.files().update(fileId=file_id, media_body=media).execute()
        
        # Sync update to index
        try:
            file_meta = index_data.get(file_id)
            if file_meta:
                file_meta["size"] = str(len(content))
                file_meta["modified_time"] = datetime.utcnow().strftime("%Y-%m-%dT%H:%M:%S.000Z")
                update_index_file(user_email, file_id, file_meta, raw_content=content)
        except Exception as index_err:
            print(f"Error syncing document save to index: {index_err}")
            
        return jsonify({"success": True})
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route('/api/upload', methods=['POST'])
def upload_file():
    """Upload physical file to Google Drive and queue asynchronously for AI processing"""
    if 'file' not in request.files:
        return jsonify({"success": False, "error": "No file provided"}), 400
    
    file = request.files['file']
    if file.filename == '':
        return jsonify({"success": False, "error": "No file selected"}), 400
        
    parent_id = request.form.get('parent_id', 'root')
    
    service = get_drive_service()
    if not service:
        return jsonify({"error": "Unauthorized"}), 401
        
    try:
        file_metadata = {
            'name': file.filename
        }
        if parent_id != 'root':
            file_metadata['parents'] = [parent_id]
            
        stream = io.BytesIO(file.read())
        media = MediaIoBaseUpload(
            stream,
            mimetype=file.mimetype or 'application/octet-stream',
            resumable=True
        )
        new_file = service.files().create(body=file_metadata, media_body=media, fields='id, name, webViewLink, webContentLink').execute()
        
        # Sync update to index
        try:
            user_email = get_session_email_or_fetch(service)
            parent_path = "/"
            if parent_id != 'root':
                index_data = load_json_file(f"drive_index_{sanitize_filename(user_email)}.json", {})
                parent_meta = index_data.get(parent_id)
                if parent_meta:
                    parent_path = parent_meta["path"]
            
            file_path = (parent_path + "/" + file.filename) if parent_path != "/" else ("/" + file.filename)
            if not file_path.startswith("/"):
                file_path = "/" + file_path
 
            stream.seek(0, io.SEEK_END)
            file_size_bytes = stream.tell()
            stream.seek(0)
 
            file_meta = {
                "id": new_file["id"],
                "name": file.filename,
                "mime_type": file.mimetype or 'application/octet-stream',
                "path": file_path,
                "parent_id": parent_id,
                "size": str(file_size_bytes),
                "modified_time": datetime.utcnow().strftime("%Y-%m-%dT%H:%M:%S.000Z"),
                "starred": False,
                "owner": "me",
                "is_folder": False,
                "webViewLink": new_file.get("webViewLink"),
                "webContentLink": new_file.get("webContentLink"),
                "embedding_status": "Pending"
            }
            
            # Save basic metadata to index immediately
            update_index_file_metadata_only(user_email, new_file["id"], file_meta)
            
            # Queue background indexing job
            if is_supported_for_extraction(file_meta["mime_type"], file_meta["name"]) and file_size_bytes < 10 * 1024 * 1024:
                nvidia_key = session.get('nvidia_api_key') or os.environ.get('NVIDIA_API_KEY') or DEFAULT_NVIDIA_API_KEY
                keys_dict = {'nvidia_api_key': nvidia_key}
                creds_dict = credentials_to_dict(get_credentials())
                
                AI_PROCESSING_QUEUE.put({
                    "user_email": user_email,
                    "file_id": new_file["id"],
                    "file_name": file.filename,
                    "mime_type": file_meta["mime_type"],
                    "credentials_dict": creds_dict,
                    "keys_dict": keys_dict
                })
            else:
                update_file_status(user_email, new_file["id"], "Not Supported")
                
        except Exception as index_err:
            print(f"Error syncing file upload to index: {index_err}")
            
        return jsonify({"success": True, "file": new_file})
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route('/api/toggle-star', methods=['POST'])
def toggle_star():
    """Toggle star status on Google Drive"""
    data = request.get_json(silent=True) or {}
    file_id = data.get('id')

    if not file_id:
        return jsonify({"error": "File id is required"}), 400

    service = get_drive_service()
    if not service:
        return jsonify({"error": "Unauthorized"}), 401
        
    try:
        user_email = get_session_email_or_fetch(service)
        if not user_email:
            return jsonify({"error": "Unauthorized"}), 401
        sanitized_email = sanitize_filename(user_email)
        index_data = load_json_file(f"drive_index_{sanitized_email}.json", {})
        has_access, file_meta = verify_file_ownership(service, user_email, file_id, index_data)
        if not has_access:
            return jsonify({"error": "Forbidden - File ownership verification failed"}), 403
            
        # Get existing state
        f = service.files().get(fileId=file_id, fields='starred').execute()
        new_state = not f.get('starred', False)
        
        service.files().update(fileId=file_id, body={'starred': new_state}).execute()
        
        try:
            file_meta = index_data.get(file_id)
            if file_meta:
                file_meta["starred"] = new_state
                update_index_file(user_email, file_id, file_meta)
        except Exception as index_err:
            print(f"Error syncing toggle star to index: {index_err}")
            
        return jsonify({"success": True})
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route('/api/delete', methods=['POST'])
def delete_file():
    """Soft delete to trash. If already in trash, delete permanently."""
    data = request.get_json(silent=True) or {}
    file_id = data.get('id')

    if not file_id:
        return jsonify({"error": "File id is required"}), 400

    service = get_drive_service()
    if not service:
        return jsonify({"error": "Unauthorized"}), 401
        
    try:
        user_email = get_session_email_or_fetch(service)
        if not user_email:
            return jsonify({"error": "Unauthorized"}), 401
        sanitized_email = sanitize_filename(user_email)
        index_data = load_json_file(f"drive_index_{sanitized_email}.json", {})
        has_access, file_meta = verify_file_ownership(service, user_email, file_id, index_data)
        if not has_access:
            return jsonify({"error": "Forbidden - File ownership verification failed"}), 403
            
        f = service.files().get(fileId=file_id, fields='trashed').execute()
        if f.get('trashed', False):
            # Delete permanently
            service.files().delete(fileId=file_id).execute()
            remove_from_index(user_email, file_id)
        else:
            # Move to trash
            service.files().update(fileId=file_id, body={'trashed': True}).execute()
            remove_from_index(user_email, file_id)
            
        return jsonify({"success": True})
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route('/api/restore', methods=['POST'])
def restore_file():
    """Restore file from Google Drive trash"""
    data = request.get_json(silent=True) or {}
    file_id = data.get('id')

    if not file_id:
        return jsonify({"error": "File id is required"}), 400

    service = get_drive_service()
    if not service:
        return jsonify({"error": "Unauthorized"}), 401
        
    try:
        user_email = get_session_email_or_fetch(service)
        if not user_email:
            return jsonify({"error": "Unauthorized"}), 401
        sanitized_email = sanitize_filename(user_email)
        index_data = load_json_file(f"drive_index_{sanitized_email}.json", {})
        has_access, file_meta = verify_file_ownership(service, user_email, file_id, index_data)
        if not has_access:
            return jsonify({"error": "Forbidden - File ownership verification failed"}), 403
            
        service.files().update(fileId=file_id, body={'trashed': False}).execute()
        
        try:
            refetched = service.files().get(
                fileId=file_id,
                fields="id, name, mimeType, size, modifiedTime, starred, parents, owners, webViewLink, webContentLink"
            ).execute()
            
            parent_id = refetched.get('parents', ['root'])[0]
            parent_path = "/"
            if parent_id != 'root':
                parent_meta = index_data.get(parent_id)
                if parent_meta:
                    parent_path = parent_meta["path"]
                    
            file_path = (parent_path + "/" + refetched["name"]) if parent_path != "/" else ("/" + refetched["name"])
            if not file_path.startswith("/"):
                file_path = "/" + file_path
                
            owner = 'me'
            if refetched.get('owners'):
                owner = refetched.get('owners')[0].get('displayName', 'me')
                
            file_metadata = {
                "id": file_id,
                "name": refetched["name"],
                "mime_type": refetched["mimeType"],
                "path": file_path,
                "parent_id": parent_id,
                "size": refetched.get('size'),
                "modified_time": refetched.get('modifiedTime'),
                "starred": refetched.get('starred', False),
                "owner": owner,
                "is_folder": refetched["mimeType"] == 'application/vnd.google-apps.folder',
                "webViewLink": refetched.get('webViewLink'),
                "webContentLink": refetched.get('webContentLink')
            }
            raw_content = None
            if not file_metadata["is_folder"] and is_supported_for_extraction(file_metadata["mime_type"], file_metadata["name"]):
                raw_content = fetch_file_text_content(service, file_id)
                
            update_index_file(user_email, file_id, file_metadata, raw_content=raw_content)
        except Exception as index_err:
            print(f"Error syncing restore file to index: {index_err}")
            
        return jsonify({"success": True})
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route('/api/empty-trash', methods=['POST'])
def empty_trash():
    """Empty Google Drive trash"""
    service = get_drive_service()
    if not service:
        return jsonify({"error": "Unauthorized"}), 401
        
    try:
        service.files().emptyTrash().execute()
        return jsonify({"success": True})
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route('/api/rename', methods=['POST'])
def rename_file():
    """Rename a file or folder on Google Drive"""
    data = request.get_json(silent=True) or {}
    file_id = data.get('id')
    new_name = (data.get('new_name') or '').strip()

    if not file_id or not new_name:
        return jsonify({"error": "File id and new name are required"}), 400

    service = get_drive_service()
    if not service:
        return jsonify({"error": "Unauthorized"}), 401
        
    try:
        user_email = get_session_email_or_fetch(service)
        if not user_email:
            return jsonify({"error": "Unauthorized"}), 401
        sanitized_email = sanitize_filename(user_email)
        index_data = load_json_file(f"drive_index_{sanitized_email}.json", {})
        has_access, file_meta = verify_file_ownership(service, user_email, file_id, index_data)
        if not has_access:
            return jsonify({"error": "Forbidden - File ownership verification failed"}), 403
            
        service.files().update(fileId=file_id, body={'name': new_name}).execute()
        
        try:
            update_cached_paths_after_rename(user_email, file_id, None, new_name)
        except Exception as index_err:
            print(f"Error syncing rename to index: {index_err}")
            
        return jsonify({"success": True})
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route('/api/file-content/<file_id>')
def get_file_content(file_id):
    """Fetch raw Google Drive file media or plain text document preview data"""
    service = get_drive_service()
    if not service:
        return jsonify({"error": "Unauthorized"}), 401
        
    try:
        user_email = get_session_email_or_fetch(service)
        if not user_email:
            return jsonify({"error": "Unauthorized"}), 401
            
        sanitized_email = sanitize_filename(user_email)
        index_data = load_json_file(f"drive_index_{sanitized_email}.json", {})
        has_access, file_meta = verify_file_ownership(service, user_email, file_id, index_data)
        if not has_access:
            return jsonify({"error": "Forbidden - File ownership verification failed"}), 403
            
        # Retrieve metadata
        meta = service.files().get(fileId=file_id, fields='name, mimeType').execute()
        mime = meta.get('mimeType', 'application/octet-stream')
        
        # Store in session as active file context
        session['selected_file_id'] = file_id
        session['selected_file_name'] = meta.get('name', 'Document')
        session['selected_file_type'] = mime
        session.modified = True
        
        # Check if it's a native Google Doc, which must be exported
        if mime == 'application/vnd.google-apps.document':
            request_media = service.files().export_media(fileId=file_id, mimeType='text/plain')
            fh = io.BytesIO()
            downloader = MediaIoBaseDownload(fh, request_media)
            done = False
            while not done:
                _, done = downloader.next_chunk()
            fh.seek(0)
            return jsonify({
                "name": meta.get('name'),
                "content": fh.read().decode('utf-8', errors='ignore'),
                "is_virtual": True
            })

        # Fetch ordinary file stream
        request_media = service.files().get_media(fileId=file_id)
        fh = io.BytesIO()
        downloader = MediaIoBaseDownload(fh, request_media)
        done = False
        while not done:
            _, done = downloader.next_chunk()
        fh.seek(0)

        # Handle rendering textual contents inside app inline code editor
        name_lower = meta.get('name', '').lower()
        text_extensions = {
            '.txt', '.md', '.json', '.js', '.html', '.css', '.py', '.ts', '.c', '.cpp', 
            '.ipynb', '.csv', '.xml', '.yaml', '.yml', '.sh', '.bat', '.cmd', '.ps1', 
            '.sql', '.log', '.env', '.toml', '.rst', '.java', '.go', '.rs', '.swift', 
            '.kt', '.cs', '.h', '.hpp', '.properties', '.jsonld', '.less', '.sass', '.scss'
        }
        is_text = (
            'text/' in mime or 
            mime in ('application/json', 'application/x-ipynb+json', 'application/javascript') or 
            any(name_lower.endswith(ext) for ext in text_extensions)
        )
        if is_text:
            return jsonify({
                "name": meta.get('name'),
                "content": fh.read().decode('utf-8', errors='ignore'),
                "is_virtual": True
            })

        # Return standard response file attachment stream for image previewing / downloading
        return send_file(
            fh,
            mimetype=mime,
            as_attachment=False,
            download_name=meta.get('name')
        )
    except Exception as e:
        return jsonify({"error": str(e)}), 500



def save_document_analysis(user_email, file_id, key, value):
    """
    Save generated analysis field (like summary, keywords, entities, etc.) into file metadata.
    """
    sanitized_email = sanitize_filename(user_email)
    index_file = f"drive_index_{sanitized_email}.json"
    index_data = load_json_file(index_file, {})
    if file_id in index_data:
        index_data[file_id][key] = value
        save_json_file(index_file, index_data)

@app.route('/api/file-metadata/<file_id>')
def get_file_metadata(file_id):
    """Fetch structured metadata, cached summaries, and related documents for a file"""
    service = get_drive_service()
    if not service:
        return jsonify({"error": "Unauthorized"}), 401
        
    try:
        user_email = get_session_email_or_fetch(service)
        if not user_email:
            return jsonify({"error": "Unauthorized"}), 401
        sanitized_email = sanitize_filename(user_email)
        
        # Load index
        index_file = f"drive_index_{sanitized_email}.json"
        index_data = load_json_file(index_file, {})
        
        has_access, file_meta = verify_file_ownership(service, user_email, file_id, index_data)
        if not has_access:
            return jsonify({"error": "Forbidden - File ownership verification failed"}), 403
            
        file_meta = index_data.get(file_id, {})
        
        # Store in session as active file context
        if file_meta:
            session['selected_file_id'] = file_id
            session['selected_file_name'] = file_meta.get('name', 'Document')
            session['selected_file_type'] = file_meta.get('mime_type') or file_meta.get('mimeType', '')
            session.modified = True
        
        # If not in index, query Drive API for basic info
        if not file_meta:
            try:
                drive_meta = service.files().get(fileId=file_id, fields='name, mimeType, size, modifiedTime').execute()
                file_meta = {
                    "id": file_id,
                    "name": drive_meta.get("name", "Document"),
                    "mime_type": drive_meta.get("mimeType", ""),
                    "size": drive_meta.get("size", "0"),
                    "modified_time": drive_meta.get("modifiedTime", ""),
                    "is_folder": drive_meta.get("mimeType") == 'application/vnd.google-apps.folder'
                }
            except Exception:
                return jsonify({"error": "File not found"}), 404
                
        # Word count & Chunks count from content cache
        content_file = f"drive_content_cache_{sanitized_email}.json"
        content_data = load_json_file(content_file, {})
        
        word_count = 0
        chunks_count = 0
        embedding_status = "Not Indexed"
        
        # Count chunks in SQLite DB
        db_path = embeddings_db.get_db_path(user_email)
        if os.path.exists(db_path):
            try:
                conn = sqlite3.connect(db_path)
                cursor = conn.cursor()
                cursor.execute("SELECT COUNT(*) FROM chunks WHERE file_id = ?", (file_id,))
                chunks_count = cursor.fetchone()[0]
                conn.close()
                if chunks_count > 0:
                    embedding_status = "Indexed"
            except Exception as e:
                print(f"Error counting chunks: {e}")
                
        # Get word count from cache
        if file_id in content_data:
            chunks = content_data[file_id].get("chunks", [])
            text = "\n\n".join(chunks)
            word_count = len(text.split())
            if chunks_count == 0:
                chunks_count = len(chunks)
        
        # Format size
        size_bytes = int(file_meta.get("size") or 0)
        if size_bytes > 0:
            if size_bytes < 1024:
                size_str = f"{size_bytes} B"
            elif size_bytes < 1024 * 1024:
                size_str = f"{round(size_bytes / 1024, 1)} KB"
            else:
                size_str = f"{round(size_bytes / (1024 * 1024), 1)} MB"
        else:
            size_str = "--"
            
        # Reading time (words / 200)
        reading_time = max(1, round(word_count / 200)) if word_count > 0 else 0
        
        # Build metadata dictionary
        result = {
            "id": file_id,
            "name": file_meta.get("name", "Document"),
            "mime_type": file_meta.get("mime_type") or file_meta.get("mimeType", ""),
            "size_str": size_str,
            "size_bytes": size_bytes,
            "word_count": word_count,
            "reading_time": reading_time,
            "modified_time": file_meta.get("modified_time") or file_meta.get("modifiedTime", ""),
            "last_indexed": file_meta.get("last_indexed") or file_meta.get("modified_time") or "--",
            "chunks_count": chunks_count,
            "embedding_status": embedding_status,
            # Document memory / cached items
            "summary": file_meta.get("summary", ""),
            "keywords": file_meta.get("keywords", ""),
            "entities": file_meta.get("entities", ""),
            "topics": file_meta.get("topics", ""),
            "important_points": file_meta.get("important_points", ""),
            "faq": file_meta.get("faq", "")
        }
        return jsonify(result)
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route('/api/storage-info')
def get_storage_info():
    """Fetch user's actual Google Drive storage capacity and usage details"""
    service = get_drive_service()
    if not service:
        return jsonify({"error": "Unauthorized"}), 401
        
    try:
        about = service.about().get(fields="storageQuota").execute()
        quota = about.get('storageQuota', {})
        
        limit_val = quota.get('limit')
        usage_val = quota.get('usage')
        
        # Convert to int safely, with fallback if limit is missing or None
        limit = int(limit_val) if limit_val is not None else (15 * 1024 * 1024 * 1024)
        usage = int(usage_val) if usage_val is not None else 0
        
        limit_gb = limit / (1024 * 1024 * 1024)
        usage_gb = usage / (1024 * 1024 * 1024)
        
        user_email = session.get('user_email')
        if not user_email:
            try:
                user_email = get_user_email(service)
                session['user_email'] = user_email
            except Exception:
                user_email = ""
            
        return jsonify({
            "used_bytes": usage,
            "used_gb": round(usage_gb, 2),
            "total_gb": round(limit_gb, 1),
            "percentage": round((usage / limit) * 100, 1) if limit else 0,
            "email": user_email or ""
        })
    except Exception as e:
        logger.exception("Error in get_storage_info")
        return jsonify({"error": str(e)}), 500

# ========= DRIVE AI ASSISTANT ENDPOINTS & TOOLS =========

def search_drive_item(service, name, mime_type=None, parent_id=None):
    """Search for an item (file or folder) in Google Drive by name"""
    try:
        safe_name = (name or "").replace("\\", "\\\\").replace("'", "\\'")
        q = f"name = '{safe_name}' and trashed = false"
        if mime_type:
            q += f" and mimeType = '{mime_type}'"
        if parent_id and parent_id != 'root':
            q += f" and '{parent_id}' in parents"
        
        results = service.files().list(q=q, fields="files(id, name, mimeType, parents)").execute()
        files = results.get('files', [])
        return files[0] if files else None
    except Exception as e:
        print(f"Error searching drive item '{name}': {e}")
        return None

def fetch_file_text_content(service, file_id):
    """Fetch text content of a Google Drive file for tool analysis.

    Handles plain text, Google Docs, PDFs, Word/PowerPoint documents, and
    images (via OCR), mirroring the extraction logic used by the background
    indexing pipeline so live (not-yet-indexed) files are read correctly.
    """
    try:
        meta = service.files().get(fileId=file_id, fields='name, mimeType').execute()
        mime = meta.get('mimeType', '')
        name = meta.get('name', '')

        if mime == 'application/vnd.google-apps.folder':
            return f"Folder: {name}"

        if mime == 'application/vnd.google-apps.document':
            # Export Google Doc as text
            req = service.files().export_media(fileId=file_id, mimeType='text/plain')
            fh = io.BytesIO()
            downloader = MediaIoBaseDownload(fh, req)
            done = False
            while not done:
                _, done = downloader.next_chunk()
            fh.seek(0)
            return fh.read().decode('utf-8', errors='ignore')

        # Download ordinary file
        req = service.files().get_media(fileId=file_id)
        fh = io.BytesIO()
        downloader = MediaIoBaseDownload(fh, req)
        done = False
        while not done:
            _, done = downloader.next_chunk()
        fh.seek(0)
        file_bytes = fh.read()

        return extract_text_from_bytes(file_bytes, mime, name)
    except Exception as e:
        print(f"Error reading file content for {file_id}: {e}")
        return ""

# --- True Google Drive Agent Indexing & Retrieval Engine ---

def sanitize_filename(name):
    if not name or name == 'default_user':
        raise ValueError("Invalid user email/name for file operations")
    return re.sub(r'[^a-zA-Z0-9_\-]', '_', name)

def get_user_data_path(filename):
    os.makedirs('embedded_files', exist_ok=True)
    return os.path.join('embedded_files', filename)

def load_json_file(filepath, default):
    if filepath.startswith("drive_index_") or filepath.startswith("drive_content_cache_"):
        filepath = get_user_data_path(filepath)
    if not os.path.exists(filepath):
        return default
    with index_lock:
        try:
            with open(filepath, 'r', encoding='utf-8') as f:
                return json.load(f)
        except Exception as e:
            print(f"Error reading {filepath}: {e}")
            return default

def save_json_file(filepath, data):
    if filepath.startswith("drive_index_") or filepath.startswith("drive_content_cache_"):
        filepath = get_user_data_path(filepath)
    with index_lock:
        try:
            with open(filepath, 'w', encoding='utf-8') as f:
                json.dump(data, f, ensure_ascii=False, indent=2)
        except Exception as e:
            print(f"Error writing {filepath}: {e}")

def load_user_cache(email):
    logger.info(f"Loading user cache for {email}")
    sanitized = sanitize_filename(email)
    return load_json_file(f"drive_content_cache_{sanitized}.json", {})

def load_user_index(email):
    logger.info(f"Loading user index for {email}")
    sanitized = sanitize_filename(email)
    return load_json_file(f"drive_index_{sanitized}.json", {})

def connect_user_vectors(email):
    logger.info(f"Connecting to vector database for {email}")
    embeddings_db.init_db(email)

def restore_user_state(email):
    logger.info(f"Restoring user state for {email}")
    if email not in INDEXING_STATUS:
        INDEXING_STATUS[email] = {"status": "not_started"}

def get_user_email(service):
    """Fetch user's email from Google Drive API about endpoint.
    Raises ValueError if email cannot be fetched.
    """
    try:
        about = service.about().get(fields="user").execute()
        email = about.get('user', {}).get('emailAddress')
        if email and email != 'default_user':
            return email
    except Exception as e:
        logger.error(f"Error fetching user email: {e}")
    raise ValueError("Failed to retrieve user email from Google API")

def extract_pdf_text(file_bytes):
    try:
        from pypdf import PdfReader
        reader = PdfReader(io.BytesIO(file_bytes))
        text = ""
        for page in reader.pages:
            t = page.extract_text()
            if t:
                text += t + "\n"
        return text
    except Exception as e:
        print(f"Error parsing PDF: {e}")
        return ""

def extract_docx_text(file_bytes):
    try:
        from docx import Document
        doc = Document(io.BytesIO(file_bytes))
        text = "\n".join([p.text for p in doc.paragraphs])
        for table in doc.tables:
            for row in table.rows:
                text += "\n" + " | ".join([cell.text for cell in row.cells])
        return text
    except Exception as e:
        print(f"Error parsing DOCX: {e}")
        return ""

def extract_pptx_text(file_bytes):
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(file_bytes))
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text += shape.text + "\n"
        return text
    except Exception as e:
        print(f"Error parsing PPTX: {e}")
        return ""

SUPPORTED_TEXT_MIMES = [
    'text/plain', 'text/markdown', 'text/html', 'text/css', 'text/csv', 
    'application/json', 'application/javascript', 'application/xml'
]

SUPPORTED_IMAGE_MIMES = [
    'image/png', 'image/jpeg', 'image/jpg', 'image/gif', 'image/webp'
]

def extract_xlsx_text(file_bytes):
    try:
        import zipfile
        import xml.etree.ElementTree as ET
        fh = io.BytesIO(file_bytes)
        with zipfile.ZipFile(fh) as z:
            # 1. Read shared strings
            shared_strings = []
            if 'xl/sharedStrings.xml' in z.namelist():
                with z.open('xl/sharedStrings.xml') as f:
                    tree = ET.parse(f)
                    root = tree.getroot()
                    ns = {'ns': 'http://schemas.openxmlformats.org/spreadsheetml/2006/main'}
                    for t in root.findall('.//ns:t', ns):
                        if t.text:
                            shared_strings.append(t.text)
            
            # 2. Extract cell data from sheets
            sheet_texts = []
            for name in z.namelist():
                if name.startswith('xl/worksheets/sheet') and name.endswith('.xml'):
                    with z.open(name) as f:
                        tree = ET.parse(f)
                        root = tree.getroot()
                        ns = {'ns': 'http://schemas.openxmlformats.org/spreadsheetml/2006/main'}
                        for v in root.findall('.//ns:v', ns):
                            if v.text:
                                sheet_texts.append(v.text)
                                
            # Combine shared strings and sheet numbers
            all_text = " ".join(shared_strings + sheet_texts)
            if not all_text.strip():
                return "Empty Spreadsheet"
            return all_text
    except Exception as e:
        print(f"Error parsing XLSX: {e}")
        return ""

def extract_binary_strings(file_bytes):
    import string
    printable = set(string.printable.encode('ascii'))
    current = bytearray()
    strings = []
    for b in file_bytes:
        if b in printable:
            current.append(b)
        else:
            if len(current) >= 4:
                try:
                    strings.append(current.decode('ascii', errors='ignore'))
                except Exception:
                    pass
            current = bytearray()
    if len(current) >= 4:
        strings.append(current.decode('ascii', errors='ignore'))
    return " ".join(strings)

def extract_text_from_bytes(file_bytes, mime_type, filename, keys_dict=None):
    if not file_bytes:
        return ""
        
    mime = mime_type.lower()
    name = filename.lower()
    ext = os.path.splitext(name)[1].lstrip('.')
    
    # 1. PDF
    if mime == 'application/pdf' or ext == 'pdf':
        return extract_pdf_text(file_bytes)
        
    # 2. DOCX
    if mime == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document' or ext == 'docx':
        return extract_docx_text(file_bytes)
        
    # 3. PPTX
    if mime == 'application/vnd.openxmlformats-officedocument.presentationml.presentation' or ext == 'pptx':
        return extract_pptx_text(file_bytes)
        
    # 4. XLSX
    if mime == 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet' or ext == 'xlsx':
        return extract_xlsx_text(file_bytes)
        
    # 5. IMAGES (except SVG)
    image_mimes = ['image/png', 'image/jpeg', 'image/jpg', 'image/gif', 'image/webp']
    image_exts = ['png', 'jpg', 'jpeg', 'gif', 'webp']
    if mime in image_mimes or ext in image_exts:
        if not keys_dict and has_request_context():
            nvidia_key = session.get('nvidia_api_key') or os.environ.get('NVIDIA_API_KEY') or DEFAULT_NVIDIA_API_KEY
            keys_dict = {'nvidia_api_key': nvidia_key}
        elif not keys_dict:
            nvidia_key = os.environ.get('NVIDIA_API_KEY') or DEFAULT_NVIDIA_API_KEY
            keys_dict = {'nvidia_api_key': nvidia_key}
        return run_nvidia_image_ocr(file_bytes, mime_type, keys_dict)
        
    # 6. AUDIO & VIDEO
    audio_video_mimes = ['audio/', 'video/', 'application/ogg']
    audio_video_exts = ['mp3', 'wav', 'm4a', 'aac', 'mp4', 'mov', 'avi', 'mkv', 'webm', 'ogg']
    if any(m in mime for m in audio_video_mimes) or ext in audio_video_exts:
        return f"Media File metadata: Name={filename}, MimeType={mime_type}, Size={len(file_bytes)} bytes."
        
    # 7. Text/Code/SVG/Notebooks
    text_mimes = ['text/', 'application/json', 'application/xml', 'application/javascript', 'image/svg+xml']
    text_exts = [
        'txt', 'csv', 'json', 'xml', 'md', 'log', 'svg',
        'py', 'js', 'ts', 'java', 'cpp', 'c', 'html', 'css', 'php', 'sql',
        'ipynb'
    ]
    if any(m in mime for m in text_mimes) or ext in text_exts:
        return file_bytes.decode('utf-8', errors='ignore')
        
    # 8. DOC, PPT, XLS (old binary documents)
    binary_doc_exts = ['doc', 'ppt', 'xls']
    if ext in binary_doc_exts:
        return extract_binary_strings(file_bytes)
        
    # 9. Fallback (best effort decode, else return metadata)
    try:
        return file_bytes.decode('utf-8')
    except UnicodeDecodeError:
        return f"Binary File metadata: Name={filename}, MimeType={mime_type}, Size={len(file_bytes)} bytes."

def is_supported_for_extraction(mime_type, filename):
    return mime_type != 'application/vnd.google-apps.folder'

def chunk_text(text, chunk_size=1000, overlap=200):
    chunks = []
    if not text:
        return chunks
    
    start = 0
    while start < len(text):
        end = start + chunk_size
        chunks.append(text[start:end])
        start += chunk_size - overlap
        if end >= len(text):
            break
    return chunks

def update_file_status(user_email, file_id, status):
    """Update only the embedding status field of a file"""
    sanitized_email = sanitize_filename(user_email)
    index_file = f"drive_index_{sanitized_email}.json"
    index_data = load_json_file(index_file, {})
    if file_id in index_data:
        index_data[file_id]["embedding_status"] = status
        save_json_file(index_file, index_data)

def update_index_file_metadata_only(user_email, file_id, file_metadata):
    """Write metadata changes without triggering asynchronous background queue pipeline"""
    sanitized_email = sanitize_filename(user_email)
    index_file = f"drive_index_{sanitized_email}.json"
    index_data = load_json_file(index_file, {})
    index_data[file_id] = file_metadata
    save_json_file(index_file, index_data)

def generate_ai_knowledge(text, filename, keys_dict, provider):
    """Generate structured summary, keywords, entities, and points from document text"""
    word_count = len(text.split()) if text else 0
    reading_time = max(1, round(word_count / 200)) if word_count > 0 else 0
    
    summary = ""
    keywords = []
    important_points = []
    entities = []
    
    if not text:
        return {
            "summary": "Document Viewer Available. AI Analysis is still processing.",
            "keywords": [],
            "important_points": [],
            "entities": [],
            "reading_time": 0,
            "word_count": 0
        }
        
    api_key = keys_dict.get('nvidia_api_key') or os.environ.get('NVIDIA_API_KEY') or DEFAULT_NVIDIA_API_KEY
    
    if api_key:
        try:
            prompt = (
                "Analyze the document text and generate a structured JSON object with the following fields:\n"
                "- \"summary\": A 2-3 sentence overview of the document.\n"
                "- \"keywords\": A list of 5-8 relevant single-word keywords.\n"
                "- \"important_points\": A list of 4-6 key takeaways.\n"
                "- \"entities\": A list of notable names, organizations, places, or dates.\n\n"
                "Respond ONLY with a valid JSON block containing these fields. Do not include markdown code block characters like ```json. Just raw JSON text.\n\n"
                f"Document: {filename}\n"
                f"Text excerpt:\n{text[:10000]}"
            )
            
            res_text = None
            import requests
            url = "https://integrate.api.nvidia.com/v1/chat/completions"
            headers = {"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"}
            payload = {
                "model": "meta/llama-3.1-8b-instruct",
                "messages": [
                    {"role": "system", "content": "You are a document analyzer. Output JSON only."},
                    {"role": "user", "content": prompt}
                ],
                "temperature": 0.2,
                "max_tokens": 1024
            }
            res = requests.post(url, headers=headers, json=payload, timeout=20)
            if res.status_code == 200:
                res_text = res.json()['choices'][0]['message']['content']
                    
            if res_text:
                clean = res_text.strip()
                if "```" in clean:
                    match = re.search(r'```(?:json)?\s*(.*?)\s*```', clean, re.DOTALL)
                    if match:
                        clean = match.group(1).strip()
                
                data = json.loads(clean)
                summary = data.get("summary", "")
                keywords = data.get("keywords", [])
                important_points = data.get("important_points", [])
                entities = data.get("entities", [])
        except Exception as e:
            print(f"Error generating AI Knowledge card: {e}")
            
    # Heuristic fallbacks
    if not summary:
        summary = f"Automatic summary: The document '{filename}' contains text with {word_count} words."
        lines = [l.strip() for l in text.split('\n') if len(l.strip()) > 20]
        if lines:
            summary = " ".join(lines[:2])
            if len(summary) > 200:
                summary = summary[:197] + "..."
                
    if not keywords:
        words = re.findall(r'\b\w{6,}\b', text.lower())
        freq = {}
        for w in words:
            freq[w] = freq.get(w, 0) + 1
        keywords = [w[0] for w in sorted(freq.items(), key=lambda x: x[1], reverse=True)[:8]]
        
    if not important_points:
        lines = [l.strip() for l in text.split('\n') if len(l.strip()) > 30]
        important_points = lines[:4] if lines else ["Contains general text contents."]
        
    if not entities:
        found_ents = re.findall(r'\b[A-Z][a-z]+(?:\s+[A-Z][a-z]+)*\b', text)
        entities = list(set(found_ents))[:6]
        
    return {
        "summary": summary,
        "keywords": keywords,
        "important_points": important_points,
        "entities": entities,
        "reading_time": reading_time,
        "word_count": word_count
    }

def run_nvidia_image_ocr(file_bytes, mime_type, keys_dict):
    """Perform OCR on an image file utilizing the Nvidia multimodal model"""
    api_key = keys_dict.get('nvidia_api_key') or os.environ.get('NVIDIA_API_KEY') or DEFAULT_NVIDIA_API_KEY
    if not api_key:
        logger.warning("Nvidia API key is not configured for image OCR.")
        return "⚠️ Image upload succeeded, but OCR requires an Nvidia API key."
    if not file_bytes:
        logger.warning("Empty image payload received for OCR.")
        return "⚠️ OCR failed: the image file was empty or could not be downloaded."
    try:
        import base64
        import requests
        image_b64 = base64.b64encode(file_bytes).decode('utf-8')
        
        url = "https://integrate.api.nvidia.com/v1/chat/completions"
        headers = {
            "Authorization": f"Bearer {api_key}",
            "Content-Type": "application/json"
        }
        
        prompt = "Perform OCR on this image. Extract and transcribe all visible text as accurately as possible. Output ONLY the extracted text, preserving formatting and lines. Do not add any explanation or markdown wrapping."
        
        payload = {
            "model": "meta/llama-3.2-11b-vision-instruct",
            "messages": [
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "text",
                            "text": prompt
                        },
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:{mime_type};base64,{image_b64}"
                            }
                        }
                    ]
                }
            ],
            "max_tokens": 1024,
            "temperature": 0.0
        }
        
        response = requests.post(url, headers=headers, json=payload, timeout=35)
        if response.status_code == 200:
            res_json = response.json()
            if 'choices' in res_json and len(res_json['choices']) > 0:
                return res_json['choices'][0]['message']['content']
        else:
            # Fallback to nvidia/neva-22b
            logger.warning(f"meta/llama-3.2-11b-vision-instruct failed with code {response.status_code}. Retrying with nvidia/neva-22b...")
            payload["model"] = "nvidia/neva-22b"
            response = requests.post(url, headers=headers, json=payload, timeout=35)
            if response.status_code == 200:
                res_json = response.json()
                if 'choices' in res_json and len(res_json['choices']) > 0:
                    return res_json['choices'][0]['message']['content']
            
            logger.error(f"Nvidia OCR failed. Neva-22b response: {response.text}")
            return f"⚠️ OCR failed: Nvidia API returned code {response.status_code}"
            
        return ""
    except Exception as e:
        logger.exception("Nvidia image OCR error")
        return f"⚠️ OCR failed: {e}"

def process_single_file_job(user_email, file_id, file_name, mime_type, creds_dict, keys_dict):
    """Workflow: Process File -> Extract -> Chunk -> Embeddings -> Summary -> Keywords -> Save Cache"""
    sanitized_email = sanitize_filename(user_email)
    index_file = f"drive_index_{sanitized_email}.json"
    content_file = f"drive_content_cache_{sanitized_email}.json"
    
    update_file_status(user_email, file_id, "Processing")
    
    try:
        creds = google.oauth2.credentials.Credentials(**creds_dict)
        service = build('drive', 'v3', credentials=creds, static_discovery=True)
        token_email = get_user_email(service)
        if token_email != user_email:
            raise ValueError(f"Mismatched credentials for background file processing: {token_email} vs {user_email}")
        
        text = ""
        if mime_type == 'application/vnd.google-apps.document':
            request_media = service.files().export_media(fileId=file_id, mimeType='text/plain')
            fh = io.BytesIO()
            downloader = MediaIoBaseDownload(fh, request_media)
            done = False
            while not done:
                _, done = downloader.next_chunk()
            text = fh.getvalue().decode('utf-8', errors='ignore')
        else:
            request_media = service.files().get_media(fileId=file_id)
            fh = io.BytesIO()
            downloader = MediaIoBaseDownload(fh, request_media)
            done = False
            while not done:
                _, done = downloader.next_chunk()
            file_bytes = fh.getvalue()
            
            text = extract_text_from_bytes(file_bytes, mime_type, file_name, keys_dict)
        
        # Save full document text to SQLite documents table
        embeddings_db.save_document(user_email, file_id, file_name, mime_type, text)
        
        chunks = chunk_text(text)
        
        # Generate and save Embeddings
        nvidia_key = keys_dict.get('nvidia_api_key') or os.environ.get('NVIDIA_API_KEY') or DEFAULT_NVIDIA_API_KEY
        provider = 'nvidia'
        api_key = nvidia_key
        
        if api_key and chunks:
            embeddings = []
            batch_size = 32
            for i in range(0, len(chunks), batch_size):
                batch = chunks[i:i+batch_size]
                embs = embeddings_db.get_embeddings_batch(batch, api_key, provider=provider, is_query=False)
                if embs:
                    embeddings.extend(embs)
                else:
                    dim = 768 if provider == 'gemini' else 1024
                    embeddings.extend([[0.0]*dim] * len(batch))
            
            if embeddings:
                index_data = load_json_file(index_file, {})
                f_meta = index_data.get(file_id, {})
                f_path = f_meta.get("path", f"/{file_name}")
                embeddings_db.save_file_chunks(
                    user_email,
                    file_id,
                    file_name,
                    f_path,
                    chunks,
                    embeddings
                )
        
        # Generate Knowledge Cards metadata
        knowledge = generate_ai_knowledge(text, file_name, keys_dict, provider)
        
        # Save cache
        content_data = load_json_file(content_file, {})
        content_data[file_id] = {
            "file_id": file_id,
            "file_name": file_name,
            "mime_type": mime_type,
            "full_text": text,
            "summary": knowledge.get("summary", ""),
            "metadata": knowledge,
            "chunks": chunks
        }
        save_json_file(content_file, content_data)
        
        # Update metadata in index_data
        index_data = load_json_file(index_file, {})
        if file_id in index_data:
            index_data[file_id]["embedding_status"] = "Indexed"
            index_data[file_id]["summary"] = knowledge.get("summary", "")
            index_data[file_id]["keywords"] = ", ".join(knowledge.get("keywords", []))
            index_data[file_id]["entities"] = ", ".join(knowledge.get("entities", []))
            index_data[file_id]["important_points"] = "\n".join([f"- {p}" for p in knowledge.get("important_points", [])])
            index_data[file_id]["word_count"] = knowledge.get("word_count", 0)
            index_data[file_id]["reading_time"] = knowledge.get("reading_time", 0)
            index_data[file_id]["last_indexed"] = datetime.utcnow().strftime("%Y-%m-%dT%H:%M:%S.000Z")
            save_json_file(index_file, index_data)
            
    except Exception as ex:
        print(f"Error processing background job for file {file_id}: {ex}")
        traceback.print_exc()
        update_file_status(user_email, file_id, "Failed")

def bg_processing_worker():
    while True:
        try:
            job = AI_PROCESSING_QUEUE.get()
            if job is None:
                break
            
            user_email = job["user_email"]
            file_id = job["file_id"]
            file_name = job["file_name"]
            mime_type = job["mime_type"]
            creds_dict = job["credentials_dict"]
            keys_dict = job["keys_dict"]
            
            process_single_file_job(user_email, file_id, file_name, mime_type, creds_dict, keys_dict)
            AI_PROCESSING_QUEUE.task_done()
        except Exception as e:
            print(f"Error in bg_processing_worker: {e}")
            traceback.print_exc()

# Start background queue thread
threading.Thread(target=bg_processing_worker, daemon=True).start()

def run_background_indexing(credentials_dict, user_email, keys_dict):
    sanitized_email = sanitize_filename(user_email)
    index_file = get_user_data_path(f"drive_index_{sanitized_email}.json")
    content_file = get_user_data_path(f"drive_content_cache_{sanitized_email}.json")

    # Quick pre-loading check
    if os.path.exists(index_file):
        INDEXING_STATUS[user_email] = {
            "status": "completed",
            "files_indexed": 0,
            "total_files": 0,
            "syncing": True
        }
    else:
        INDEXING_STATUS[user_email] = {
            "status": "indexing",
            "files_indexed": 0,
            "total_files": 0
        }

    try:
        creds = google.oauth2.credentials.Credentials(**credentials_dict)
        service = build('drive', 'v3', credentials=creds, static_discovery=True)
        token_email = get_user_email(service)
        if token_email != user_email:
            raise ValueError(f"Mismatched credentials for background indexing: {token_email} vs {user_email}")
        
        # 1. Fetch flat list of all files
        all_files = []
        page_token = None
        while True:
            results = service.files().list(
                q="trashed = false",
                fields="nextPageToken, files(id, name, mimeType, size, modifiedTime, createdTime, starred, parents, owners, webViewLink, webContentLink)",
                pageSize=1000,
                pageToken=page_token
            ).execute()
            all_files.extend(results.get('files', []))
            page_token = results.get('nextPageToken')
            if not page_token:
                break
                
        total_files = len(all_files)
        INDEXING_STATUS[user_email]["total_files"] = total_files
        INDEXING_STATUS[user_email]["status"] = "indexing" # force to indexing to track progress
        
        # 2. Build folder path tree
        parent_map = {}
        name_map = {}
        for f in all_files:
            file_id = f.get('id')
            parent_map[file_id] = f.get('parents', [])
            name_map[file_id] = f.get('name', '')
            
        memo = {}
        def get_path(f_id):
            if f_id in memo:
                return memo[f_id]
            if not f_id or f_id not in parent_map:
                return "/"
            parents = parent_map[f_id]
            if not parents:
                return "/" + name_map.get(f_id, "")
            p_id = parents[0]
            if p_id == 'root' or p_id not in name_map:
                res_path = "/" + name_map.get(f_id, "")
            else:
                res_path = get_path(p_id) + "/" + name_map.get(f_id, "")
            memo[f_id] = res_path
            return res_path
            
        # Load existing index data for delta comparison
        old_index_data = load_json_file(index_file, {})
        old_content_data = load_json_file(content_file, {})
        
        # Initialize SQLite DB
        embeddings_db.init_db(user_email)
        
        # Identify deleted files to clean up vector DB
        current_file_ids = {f.get('id') for f in all_files if f.get('id')}
        for old_fid in list(old_index_data.keys()):
            if old_fid not in current_file_ids:
                embeddings_db.delete_file_chunks(user_email, old_fid)
                
        index_data = {}
        content_data = {}
        
        for idx, f in enumerate(all_files):
            file_id = f.get('id')
            name = f.get('name', '')
            mime_type = f.get('mimeType', '')
            size = f.get('size')
            modified_time = f.get('modifiedTime', '')
            starred = f.get('starred', False)
            
            path = get_path(file_id)
            is_folder = mime_type == 'application/vnd.google-apps.folder'
            
            owner = 'me'
            if f.get('owners'):
                owner = f.get('owners')[0].get('displayName', 'me')
                
            index_data[file_id] = {
                "id": file_id,
                "name": name,
                "mime_type": mime_type,
                "path": path,
                "parent_id": parent_map[file_id][0] if parent_map[file_id] else 'root',
                "size": size,
                "modified_time": modified_time,
                "created_time": f.get('createdTime') or modified_time,
                "starred": starred,
                "owner": owner,
                "is_folder": is_folder,
                "webViewLink": f.get('webViewLink'),
                "webContentLink": f.get('webContentLink')
            }
            
            # Check delta: if file is unchanged (name, path, modified_time same), skip indexing
            old_meta = old_index_data.get(file_id)
            is_unchanged = (
                old_meta and
                old_meta.get("modified_time") == modified_time and
                old_meta.get("path") == path and
                old_meta.get("name") == name
            )
            
            if is_unchanged:
                # Copy from old content cache
                if file_id in old_content_data:
                    content_data[file_id] = old_content_data[file_id]
                
                # Copy old status and AI knowledge metadata fields
                if old_meta:
                    index_data[file_id]["embedding_status"] = old_meta.get("embedding_status", "Indexed")
                    for k in ["summary", "keywords", "entities", "important_points", "word_count", "reading_time", "last_indexed"]:
                        if k in old_meta:
                            index_data[file_id][k] = old_meta[k]
                            
                INDEXING_STATUS[user_email]["files_indexed"] += 1
                continue
            
            # Extract content if modified/new - queue it asynchronously
            if not is_folder and is_supported_for_extraction(mime_type, name):
                index_data[file_id]["embedding_status"] = "Pending"
                creds_dict = credentials_to_dict(creds)
                
                # Copy from old content temporarily to avoid missing preview
                if file_id in old_content_data:
                    content_data[file_id] = old_content_data[file_id]
                    
                AI_PROCESSING_QUEUE.put({
                    "user_email": user_email,
                    "file_id": file_id,
                    "file_name": name,
                    "mime_type": mime_type,
                    "credentials_dict": creds_dict,
                    "keys_dict": keys_dict
                })
            else:
                index_data[file_id]["embedding_status"] = "Not Supported"
                
            INDEXING_STATUS[user_email]["files_indexed"] += 1
            
        save_json_file(index_file, index_data)
        save_json_file(content_file, content_data)
        
        INDEXING_STATUS[user_email] = {
            "status": "completed",
            "files_indexed": total_files,
            "total_files": total_files,
            "syncing": False
        }
        print(f"Drive indexing complete for {user_email}.")
    except Exception as e:
        traceback.print_exc()
        INDEXING_STATUS[user_email] = {
            "status": "failed",
            "error": str(e),
            "files_indexed": 0,
            "total_files": 0
        }

def update_index_file(user_email, file_id, file_metadata, raw_content=None):
    sanitized_email = sanitize_filename(user_email)
    index_file = f"drive_index_{sanitized_email}.json"
    content_file = f"drive_content_cache_{sanitized_email}.json"
    
    index_data = load_json_file(index_file, {})
    content_data = load_json_file(content_file, {})
    
    file_metadata["embedding_status"] = "Pending"
    index_data[file_id] = file_metadata
    
    # Store temporary cache entry
    if raw_content is not None:
        chunks = chunk_text(raw_content)
        content_data[file_id] = {
            "file_id": file_id,
            "file_name": file_metadata["name"],
            "mime_type": file_metadata.get("mime_type") or file_metadata.get("mimeType", "text/plain"),
            "full_text": raw_content,
            "summary": content_data.get(file_id, {}).get("summary", ""),
            "metadata": content_data.get(file_id, {}).get("metadata", {}),
            "chunks": chunks
        }
        save_json_file(content_file, content_data)
        
    save_json_file(index_file, index_data)
    
    # Queue for AI processing asynchronously
    from flask import has_request_context, session
    gemini_key = None
    nvidia_key = None
    creds_dict = None
    
    if has_request_context():
        gemini_key = session.get('gemini_api_key')
        nvidia_key = session.get('nvidia_api_key')
        creds = get_credentials()
        if creds:
            creds_dict = credentials_to_dict(creds)
            
    if not gemini_key:
        gemini_key = os.environ.get('GEMINI_API_KEY')
    if not nvidia_key:
        nvidia_key = os.environ.get('NVIDIA_API_KEY') or DEFAULT_NVIDIA_API_KEY
        
    if creds_dict:
        keys_dict = {'gemini_api_key': gemini_key, 'nvidia_api_key': nvidia_key}
        AI_PROCESSING_QUEUE.put({
            "user_email": user_email,
            "file_id": file_id,
            "file_name": file_metadata["name"],
            "mime_type": file_metadata.get("mime_type") or file_metadata.get("mimeType", "text/plain"),
            "credentials_dict": creds_dict,
            "keys_dict": keys_dict
        })

def remove_from_index(user_email, file_id):
    sanitized_email = sanitize_filename(user_email)
    index_file = f"drive_index_{sanitized_email}.json"
    content_file = f"drive_content_cache_{sanitized_email}.json"
    
    index_data = load_json_file(index_file, {})
    content_data = load_json_file(content_file, {})
    
    if file_id in index_data:
        del index_data[file_id]
    if file_id in content_data:
        del content_data[file_id]
        
    # Sync delete to vector store
    embeddings_db.delete_file_chunks(user_email, file_id)
        
    save_json_file(index_file, index_data)
    save_json_file(content_file, content_data)

def update_cached_paths_after_rename(user_email, folder_id, old_name, new_name):
    sanitized_email = sanitize_filename(user_email)
    index_file = f"drive_index_{sanitized_email}.json"
    content_file = f"drive_content_cache_{sanitized_email}.json"
    
    index_data = load_json_file(index_file, {})
    content_data = load_json_file(content_file, {})
    
    target_item = index_data.get(folder_id)
    if not target_item:
        return
        
    old_path = target_item["path"]
    parent_parts = old_path.split("/")[:-1]
    parent_path = "/".join(parent_parts)
    new_path = (parent_path + "/" + new_name) if parent_path else ("/" + new_name)
    if not new_path.startswith("/"):
        new_path = "/" + new_path
        
    target_item["name"] = new_name
    target_item["path"] = new_path
    
    if target_item.get("is_folder", False):
        prefix = old_path + "/"
        for f_id, f_meta in index_data.items():
            if f_meta["path"].startswith(prefix):
                suffix = f_meta["path"][len(prefix):]
                f_meta["path"] = new_path + "/" + suffix
                if f_id in content_data:
                    content_data[f_id]["path"] = f_meta["path"]
                    
        # Sync folder children paths in vector store
        embeddings_db.update_folder_paths(user_email, old_path, new_path)
    else:
        # Sync file rename in vector store
        embeddings_db.rename_file_chunks(user_email, folder_id, new_name, new_path)
                    
    if folder_id in content_data:
        content_data[folder_id]["name"] = new_name
        content_data[folder_id]["path"] = new_path
        
    save_json_file(index_file, index_data)
    save_json_file(content_file, content_data)



def compute_tfidf_relevance(query, documents):
    def tokenize(text):
        return re.findall(r'\b\w{3,}\b', text.lower())
        
    query_tokens = tokenize(query)
    if not query_tokens:
        return []
        
    df = {}
    doc_tokens_list = []
    for doc in documents:
        tokens = set(tokenize(doc['text']))
        doc_tokens_list.append(tokenize(doc['text']))
        for t in tokens:
            df[t] = df.get(t, 0) + 1
            
    num_docs = len(documents)
    scores = []
    for idx, doc in enumerate(documents):
        tokens = doc_tokens_list[idx]
        if not tokens:
            continue
        tf = {}
        for t in tokens:
            tf[t] = tf.get(t, 0) + 1
            
        score = 0.0
        for q_t in query_tokens:
            if q_t in tf:
                tf_val = tf[q_t] / len(tokens)
                idf_val = math.log((1 + num_docs) / (1 + df.get(q_t, 0))) + 1
                score += tf_val * idf_val
                
        if score > 0:
            scores.append((doc, score))
            
    scores.sort(key=lambda x: x[1], reverse=True)
    return scores

def find_matching_files(user_email, query):
    sanitized_email = sanitize_filename(user_email)
    index_file = f"drive_index_{sanitized_email}.json"
    content_file = f"drive_content_cache_{sanitized_email}.json"
    
    index_data = load_json_file(index_file, {})
    content_data = load_json_file(content_file, {})
    
    if not index_data:
        return [], []
        
    name_matches = []
    query_lower = query.lower()
    
    for f_id, meta in index_data.items():
        name = meta["name"]
        path = meta["path"]
        
        if query_lower in name.lower() or query_lower in path.lower():
            score = 1.0 if query_lower == name.lower() else 0.8
            name_matches.append((meta, score))
        else:
            ratio = difflib.SequenceMatcher(None, query_lower, name.lower()).ratio()
            if ratio > 0.4:
                name_matches.append((meta, ratio * 0.7))
                
    name_matches.sort(key=lambda x: x[1], reverse=True)
    top_name_matches = [item[0] for item in name_matches[:10]]
    
    all_chunks = []
    for f_id, content in content_data.items():
        meta = index_data.get(f_id, {})
        for idx, chunk in enumerate(content.get("chunks", [])):
            all_chunks.append({
                "file_id": f_id,
                "name": meta.get("name", content["name"]),
                "path": meta.get("path", content["path"]),
                "chunk_index": idx,
                "text": chunk
            })
            
    content_matches = []
    if all_chunks:
        tfidf_scores = compute_tfidf_relevance(query, all_chunks)
        content_matches = tfidf_scores[:15]
        
    return top_name_matches, content_matches

# --- Nvidia Llama 3.1 API Caller ---

def run_nvidia_llm(tool, doc_name, content, user_query):
    """Run Nvidia LLM (Llama 3.1 8B Instruct) for live document analysis"""
    api_key = session.get('nvidia_api_key') or os.environ.get('NVIDIA_API_KEY') or DEFAULT_NVIDIA_API_KEY
    if not api_key:
        return None
        
    try:
        import requests
        url = "https://integrate.api.nvidia.com/v1/chat/completions"
        headers = {
            "Authorization": f"Bearer {api_key}",
            "Content-Type": "application/json"
        }
        
        system_instruction = (
            "You are the Drive AI Assistant, a helpful assistant integrated into a Google Drive clone. "
            "Your objective is to execute specific tools on the user's document contents. "
            "Respond in clear, professional Markdown formatting. Do not include raw HTML tags. "
            f"You are running the tool: '{tool}' on the document '{doc_name}'. "
        )
        
        payload = {
            "model": "meta/llama-3.1-8b-instruct",
            "messages": [
                {
                    "role": "system",
                    "content": system_instruction
                },
                {
                    "role": "user",
                    "content": f"Document Content:\n\"\"\"\n{content[:15000]}\n\"\"\"\n\nUser request: {user_query}"
                }
            ],
            "temperature": 0.2,
            "max_tokens": 1024
        }
        
        response = requests.post(url, headers=headers, json=payload, timeout=15)
        if response.status_code == 200:
            res_json = response.json()
            return res_json['choices'][0]['message']['content']
        else:
            print(f"Nvidia API returned code {response.status_code}: {response.text}")
            return None
    except Exception as e:
        print(f"Error calling Nvidia API: {e}")
        return None


# run_gemini was removed as part of migrating exclusively to Nvidia APIs

# --- Chat Routing API ---

def get_folder_contents_desc(user_email, folder_id):
    sanitized_email = sanitize_filename(user_email)
    index_file = f"drive_index_{sanitized_email}.json"
    index_data = load_json_file(index_file, {})
    if not index_data:
        return "No files in current folder."
    
    contents = []
    for f_id, meta in index_data.items():
        if meta.get("parent_id") == folder_id:
            ftype = "folder" if meta.get("is_folder") else "file"
            contents.append(f"- `{meta['name']}` ({ftype}, ID: `{meta['id']}`, path: `{meta['path']}`)")
            
    if not contents:
        return "No files in current folder."
    return "\n".join(contents)

def resolve_file_by_name(user_email, query):
    """
    Search indexed file names for exact substring or fuzzy matches.
    Returns the file metadata dict if resolved, else None.
    """
    sanitized_email = sanitize_filename(user_email)
    index_file = f"drive_index_{sanitized_email}.json"
    index_data = load_json_file(index_file, {})
    if not index_data:
        return None
        
    query_lower = query.lower()
    clean_query = query_lower
    
    prefixes = [
        "summarize document", "summarize file", "summarize",
        "extract entities from", "extract entities", "entities from", "entities",
        "translate document to", "translate document", "translate",
        "simplify document", "simplify", "easier", "plain english",
        "sentiment analysis of", "sentiment of", "sentiment",
        "timeline of", "timeline", "milestones of", "milestones",
        "citation extraction from", "citations from", "citations", "citation",
        "extract text from", "extract text", "content of", "text from", "text of",
        "important points from", "important points of", "important points",
        "insights from", "insights of", "insights", "deep insights",
        "qa from", "qa of", "qa document", "qa"
    ]
    for p in prefixes:
        if clean_query.startswith(p):
            clean_query = clean_query[len(p):].strip()
            break
            
    clean_query = re.sub(r'[^\w\s\-\.]', '', clean_query).strip()
    if not clean_query:
        return None
        
    best_match = None
    best_score = 0.0
    
    for fid, meta in index_data.items():
        if meta.get("is_folder"):
            continue
        name = meta["name"]
        name_no_ext = os.path.splitext(name)[0]
        
        # Substring check
        if name_no_ext.lower() in clean_query or clean_query in name_no_ext.lower():
            score = len(name_no_ext) / max(1, len(clean_query))
            if score > best_score:
                best_score = score
                best_match = meta
                
    # Fallback to fuzzy matching
    if not best_match:
        for fid, meta in index_data.items():
            if meta.get("is_folder"):
                continue
            name = meta["name"]
            name_no_ext = os.path.splitext(name)[0]
            ratio = difflib.SequenceMatcher(None, clean_query, name_no_ext.lower()).ratio()
            if ratio > 0.45 and ratio > best_score:
                best_score = ratio
                best_match = meta
                
    return best_match


def get_document_content(user_email, service, file_id):
    """
    Guarantees document contents loading by using SQLite, JSON content cache,
    or pulling / parsing from Google Drive on-the-fly.
    """
    sanitized_email = sanitize_filename(user_email)
    index_file = f"drive_index_{sanitized_email}.json"
    index_data = load_json_file(index_file, {})
    
    # Enforce strict ownership check
    has_access, file_meta = verify_file_ownership(service, user_email, file_id, index_data)
    if not has_access:
        logger.warning(f"Access denied to file {file_id} for user {user_email}")
        return None
        
    content_file = f"drive_content_cache_{sanitized_email}.json"
    content_data = load_json_file(content_file, {})
    
    # Check if text is already present in cache
    if file_id in content_data and "full_text" in content_data[file_id]:
        return content_data[file_id]
        
    # Check if text is present in SQLite document store
    db_doc = embeddings_db.get_document(user_email, file_id)
    if db_doc and db_doc.get("full_text"):
        text = db_doc["full_text"]
        name = db_doc["file_name"]
        mime_type = db_doc["file_type"]
        chunks = chunk_text(text)
        
        # Load keys
        nvidia_key = None
        if has_request_context():
            nvidia_key = session.get('nvidia_api_key')
        if not nvidia_key:
            nvidia_key = os.environ.get('NVIDIA_API_KEY') or DEFAULT_NVIDIA_API_KEY
        keys_dict = {'nvidia_api_key': nvidia_key}
        provider = 'nvidia'
        
        knowledge = generate_ai_knowledge(text, name, keys_dict, provider)
        doc_record = {
            "file_id": file_id,
            "file_name": name,
            "mime_type": mime_type,
            "full_text": text,
            "summary": knowledge.get("summary", ""),
            "metadata": knowledge,
            "chunks": chunks
        }
        content_data[file_id] = doc_record
        save_json_file(content_file, content_data)
        
        # Sync status inside index
        index_file = f"drive_index_{sanitized_email}.json"
        index_data = load_json_file(index_file, {})
        if file_id in index_data:
            index_data[file_id]["embedding_status"] = "Indexed"
            index_data[file_id]["summary"] = knowledge.get("summary", "")
            index_data[file_id]["keywords"] = ", ".join(knowledge.get("keywords", []))
            index_data[file_id]["entities"] = ", ".join(knowledge.get("entities", []))
            index_data[file_id]["important_points"] = "\n".join([f"- {p}" for p in knowledge.get("important_points", [])])
            index_data[file_id]["word_count"] = knowledge.get("word_count", 0)
            index_data[file_id]["reading_time"] = knowledge.get("reading_time", 0)
            index_data[file_id]["last_indexed"] = datetime.utcnow().strftime("%Y-%m-%dT%H:%M:%S.000Z")
            save_json_file(index_file, index_data)
            
        return doc_record

    # Fetch details from Google Drive and compile cache entry
    index_file = f"drive_index_{sanitized_email}.json"
    index_data = load_json_file(index_file, {})
    
    meta = index_data.get(file_id, {})
    name = meta.get("name", "Document")
    mime_type = meta.get("mime_type") or meta.get("mimeType", "")
    
    if not mime_type:
        try:
            drive_meta = service.files().get(fileId=file_id, fields='name, mimeType').execute()
            name = drive_meta.get("name", "Document")
            mime_type = drive_meta.get("mimeType", "")
        except Exception:
            pass
            
    text = fetch_file_text_content(service, file_id)
    
    # Save full document text to SQLite documents table
    embeddings_db.save_document(user_email, file_id, name, mime_type, text)
    
    chunks = chunk_text(text)
    
    # Load keys
    nvidia_key = None
    if has_request_context():
        nvidia_key = session.get('nvidia_api_key')
    if not nvidia_key:
        nvidia_key = os.environ.get('NVIDIA_API_KEY') or DEFAULT_NVIDIA_API_KEY
    keys_dict = {'nvidia_api_key': nvidia_key}
    provider = 'nvidia'
    
    knowledge = generate_ai_knowledge(text, name, keys_dict, provider)
    
    doc_record = {
        "file_id": file_id,
        "file_name": name,
        "mime_type": mime_type,
        "full_text": text,
        "summary": knowledge.get("summary", ""),
        "metadata": knowledge,
        "chunks": chunks
    }
    
    content_data[file_id] = doc_record
    save_json_file(content_file, content_data)
    
    # Sync status inside index
    if file_id in index_data:
        index_data[file_id]["embedding_status"] = "Indexed"
        index_data[file_id]["summary"] = knowledge.get("summary", "")
        index_data[file_id]["keywords"] = ", ".join(knowledge.get("keywords", []))
        index_data[file_id]["entities"] = ", ".join(knowledge.get("entities", []))
        index_data[file_id]["important_points"] = "\n".join([f"- {p}" for p in knowledge.get("important_points", [])])
        index_data[file_id]["word_count"] = knowledge.get("word_count", 0)
        index_data[file_id]["reading_time"] = knowledge.get("reading_time", 0)
        index_data[file_id]["last_indexed"] = datetime.utcnow().strftime("%Y-%m-%dT%H:%M:%S.000Z")
        save_json_file(index_file, index_data)
        
    return doc_record

def get_document_text(user_email, service, file_id):
    """
    Retrieve full extracted text of a document from content cache or Drive.
    """
    doc = get_document_content(user_email, service, file_id)
    return doc.get("full_text", "")


def save_document_summary(user_email, file_id, summary):
    """
    Save generated summary into file metadata for future instant access.
    """
    sanitized_email = sanitize_filename(user_email)
    index_file = f"drive_index_{sanitized_email}.json"
    index_data = load_json_file(index_file, {})
    if file_id in index_data:
        index_data[file_id]["summary"] = summary
        save_json_file(index_file, index_data)


def clean_assistant_response(text):
    """
    Strip internal JSON markdown code blocks from response so users only see the final result.
    """
    if not text:
        return ""
    cleaned = re.sub(r'```json\s*.*?\s*```', '', text, flags=re.DOTALL).strip()
    if not cleaned:
        cleaned = "Action executed successfully."
    return cleaned


def get_file_records(user_email, service, file_ids):
    """
    Construct file records containing metadata, full extracted text, and summaries using get_document_content.
    """
    records = []
    for fid in file_ids:
        doc = get_document_content(user_email, service, fid)
        if doc:
            records.append({
                "id": fid,
                "name": doc.get("file_name", "Document"),
                "text": doc.get("full_text", ""),
                "word_count": doc.get("metadata", {}).get("word_count", 0) or len(doc.get("full_text", "").split()),
                "summary": doc.get("summary", "")
            })
    return records


def classify_intent(message, api_key, provider):
    """
    Classify the message intent using heuristics or LLM prompt.
    """
    msg_lower = message.lower()
    detected_intent = None
    detected_cat = None
    
    # 1. Fallback Heuristics first to check simple keyword matches
    # Drive Operations
    if re.search(r'\bcreate folder\b', msg_lower) or msg_lower.startswith("create folder"):
        detected_intent = "create folder"
        detected_cat = "Drive Operations"
    elif re.search(r'\b(create doc|create document|create file)\b', msg_lower):
        detected_intent = "create file"
        detected_cat = "Drive Operations"
    elif re.search(r'\b(delete|remove)\b', msg_lower):
        detected_intent = "delete"
        detected_cat = "Drive Operations"
    elif re.search(r'\brename\b', msg_lower):
        detected_intent = "rename"
        detected_cat = "Drive Operations"
    elif re.search(r'\bmove\b', msg_lower):
        detected_intent = "move"
        detected_cat = "Drive Operations"
    elif re.search(r'\bcopy\b', msg_lower):
        detected_intent = "copy"
        detected_cat = "Drive Operations"
    elif re.search(r'\bupload\b', msg_lower):
        detected_intent = "upload"
        detected_cat = "Drive Operations"
    elif re.search(r'\bdownload\b', msg_lower):
        detected_intent = "download"
        detected_cat = "Drive Operations"
    elif re.search(r'\bunstar\b', msg_lower):
        detected_intent = "unstar"
        detected_cat = "Drive Operations"
    elif re.search(r'\bstar\b', msg_lower):
        detected_intent = "star"
        detected_cat = "Drive Operations"
    elif re.search(r'\bshare\b', msg_lower):
        detected_intent = "share"
        detected_cat = "Drive Operations"
        
    # Document Analysis
    elif re.search(r'\b(document summary|summarize|summary)\b', msg_lower):
        detected_intent = "summarize"
        detected_cat = "Document Analysis"
    elif re.search(r'\b(analyze|analysis)\b', msg_lower):
        detected_intent = "analyze"
        detected_cat = "Document Analysis"
    elif re.search(r'\b(extract text|extract content|text of|content of)\b', msg_lower):
        detected_intent = "extract text"
        detected_cat = "Document Analysis"
    elif re.search(r'\b(extract entities|entity|entities|ner)\b', msg_lower):
        detected_intent = "extract entities"
        detected_cat = "Document Analysis"
    elif re.search(r'\b(insight|insights)\b', msg_lower):
        detected_intent = "insights"
        detected_cat = "Document Analysis"
    elif re.search(r'\b(compare|comparator|similarity)\b', msg_lower):
        detected_intent = "compare"
        detected_cat = "Document Analysis"
    elif re.search(r'\b(simplify|simple|easier)\b', msg_lower):
        detected_intent = "simplify"
        detected_cat = "Document Analysis"
    elif re.search(r'\btranslate\b', msg_lower):
        detected_intent = "translate"
        detected_cat = "Document Analysis"
    elif re.search(r'\b(timeline|chronology|milestones?)\b', msg_lower):
        detected_intent = "timeline"
        detected_cat = "Document Analysis"
    elif re.search(r'\b(tone|sentiments?|emotions?)\b', msg_lower):
        detected_intent = "sentiment"
        detected_cat = "Document Analysis"
    elif re.search(r'\b(citations?|references?|bibliography)\b', msg_lower):
        detected_intent = "citation"
        detected_cat = "Document Analysis"
    elif re.search(r'\b(takeaways?|key points?|important points?)\b', msg_lower):
        detected_intent = "takeaways"
        detected_cat = "Document Analysis"

    # Document Q&A
    elif "what does this file say" in msg_lower:
        detected_intent = "what does this file say"
        detected_cat = "Document Q&A"
    elif "explain from document" in msg_lower:
        detected_intent = "explain from document"
        detected_cat = "Document Q&A"
    elif "answer from file" in msg_lower:
        detected_intent = "answer from file"
        detected_cat = "Document Q&A"
    elif "topic in document" in msg_lower:
        detected_intent = "topic in document"
        detected_cat = "Document Q&A"
    elif any(term in msg_lower for term in ["phone", "email", "contact", "skill", "marks", "score", "date", "deadline", "gpa", "experience", "salary", "address", "resume", "cv", "chapter", "page", "section", "table", "figure"]):
        detected_intent = "qa_document"
        detected_cat = "Document Q&A"

    # General AI Chat
    elif re.search(r'\b(hello|hi|greet|hey)\b', msg_lower):
        detected_intent = "hello"
        detected_cat = "General AI Chat"
    elif re.search(r'\b(what is|what\'s)\b', msg_lower):
        detected_intent = "what is"
        detected_cat = "General AI Chat"
    elif re.search(r'\b(who is|who\'s)\b', msg_lower):
        detected_intent = "who is"
        detected_cat = "General AI Chat"
    elif re.search(r'\bhow\b', msg_lower):
        detected_intent = "how"
        detected_cat = "General AI Chat"
    elif re.search(r'\bwhy\b', msg_lower):
        detected_intent = "why"
        detected_cat = "General AI Chat"
    elif re.search(r'\bhelp\b', msg_lower):
        detected_intent = "help"
        detected_cat = "General AI Chat"
    elif re.search(r'\bpython\b', msg_lower):
        detected_intent = "python"
        detected_cat = "General AI Chat"
    elif re.search(r'\bc language\b', msg_lower):
        detected_intent = "c language"
        detected_cat = "General AI Chat"
    elif re.search(r'\bmachine learning\b', msg_lower):
        detected_intent = "machine learning"
        detected_cat = "General AI Chat"
    elif re.search(r'\bexplain\b', msg_lower):
        detected_intent = "explain"
        detected_cat = "General AI Chat"
        
    if detected_intent and detected_cat:
        return {"category": detected_cat, "intent": detected_intent, "confidence": 1.0}

    # 2. If no heuristic, call LLM to classify
    if api_key:
        CLASSIFIER_SYSTEM = """You are an intent classification system for a Drive AI Assistant.
Analyze the user's message and classify it into exactly one category and one intent from this list:

## Drive Operations:
- create folder
- create file
- delete
- rename
- move
- copy
- upload
- download
- share
- star
- unstar

## Document Analysis:
- summarize
- analyze
- extract text
- extract entities
- insights
- compare
- document summary

## Document Q&A:
- what does this file say
- explain from document
- answer from file
- topic in document
- phone number
- email
- skills
- marks
- dates
- qa_document (for general questions about document content, e.g. candidates, details, tables, specifications)

## General AI Chat:
- what is
- who is
- how
- why
- hello
- help
- python
- c language
- machine learning

Respond ONLY with a valid JSON block of this format:
{
  "category": "Category Name",
  "intent": "intent name",
  "confidence": 0.95
}"""
        try:
            raw = run_nvidia_llm("classifier", "intent", "Classifier input", message)
            if raw:
                clean = raw.strip()
                if "```" in clean:
                    clean = clean.split("```")[1]
                    if clean.startswith("json"):
                        clean = clean[4:]
                data = json.loads(clean.strip())
                return {
                    "category": data.get("category", "General AI Chat"),
                    "intent": data.get("intent", "help"),
                    "confidence": float(data.get("confidence", 0.5))
                }
        except Exception as e:
            print(f"Classifier LLM error: {e}")
            
    return {"category": "General AI Chat", "intent": "help", "confidence": 0.5}


@app.route('/api/ai/chat', methods=['POST'])
def ai_chat():
    """Chatbot assistant processing endpoint"""
    try:
        data = request.get_json(silent=True) or {}
        message = str(data.get("message", "")).strip()
        msg_lower = message.lower()
        
        logger.info(f"Incoming message: {message}")
        
        selected_file_id = data.get('selected_file_id') or session.get('selected_file_id')
        selected_file_ids = data.get('selected_file_ids') or []
        attached_files = data.get('attached_files') or []
        current_folder_id = data.get('folder_id', 'root')
        
        logger.info(f"Selected file: {selected_file_id}")
        
        service = get_drive_service()
        if not service:
            return jsonify({"error": "Unauthorized"}), 401
     
        if not message:
            return jsonify({"error": "Empty message"}), 400
     
        # Ensure chat history is initialized
        if 'chat_history' not in session:
            session['chat_history'] = []
     
        # Limit session chat history (max 8 messages)
        history = session['chat_history']
        if len(history) > 8:
            history = history[-8:]
     
        def get_updated_history(role, text):
            history.append({
                "role": role,
                "content": text,
                "timestamp": datetime.now().strftime("%H:%M")
            })
            session['chat_history'] = history
            session.modified = True
            return history
     
        # Append user message
        get_updated_history("user", message)
     
        user_email = get_session_email_or_fetch(service)
        sanitized_email = sanitize_filename(user_email)
     
        # 1. API Key setting command
        if msg_lower.startswith('set key ') or msg_lower.startswith('set api key '):
            parts = message.split(' ')
            key = parts[-1].strip()
            if len(key) > 10:
                session['nvidia_api_key'] = key
                session.modified = True
                msg = "🔑 NVIDIA API Key set successfully! The Drive AI Assistant will now use NVIDIA models for all operations."
                return jsonify({"response": msg, "history": get_updated_history("assistant", msg)})
            else:
                msg = "❌ Invalid API Key format."
                return jsonify({"response": msg, "history": get_updated_history("assistant", msg)})
     

        # Get active key and provider
        nvidia_key = session.get('nvidia_api_key') or os.environ.get('NVIDIA_API_KEY') or DEFAULT_NVIDIA_API_KEY
        provider = 'nvidia'
        api_key = nvidia_key
     
        # Run Intent Classification Layer
        classification = classify_intent(message, api_key, provider)
        category = classification["category"]
        intent = classification["intent"]
        confidence = classification["confidence"]
        
        logger.info(f"Intent: {intent}")
     
        # Unify selected_file_id and selected_file_ids list
        if not selected_file_ids and selected_file_id:
            selected_file_ids = [selected_file_id]
     
        # Route 1: Drive Operations (Manual Guidance Only)
        if category == "Drive Operations":
            logger.info("Executing manual-only Drive Operations routing path")
            
            instructions = {
                "delete": "🗑️ **Delete File/Folder**: To delete a file or folder manually, right-click on the item in the list and select **Delete** from the context menu.",
                "move": "📂 **Move File/Folder**: To move a file or folder manually, right-click on the item, select **Move**, and select the destination folder.",
                "copy": "📋 **Copy File**: To copy a file manually, right-click on the file and select **Make a Copy**.",
                "upload": "📤 **Upload File**: To upload a file manually, click the **Upload File** button at the top-left of the sidebar, or simply drag and drop the file into the files grid/list.",
                "download": "📥 **Download File**: To download a file manually, right-click on it and select **Download**, or click the **Download** button inside the file preview.",
                "star": "⭐ **Star File**: To star a file manually, right-click on the file and select **Star**.",
                "unstar": "☆ **Unstar File**: To remove a star from a file manually, right-click on the file and select **Unstar**.",
                "share": "🔗 **Share File**: To share a file manually (allowing anyone with the link to read it), right-click on the file and select **Share**."
            }
            
            resp = instructions.get(intent, "⚠️ **Drive Operations**: The AI Assistant cannot modify your Drive files directly. Please perform file actions (creating, deleting, renaming, moving, copying, etc.) manually using the buttons in the sidebar or by right-clicking on any file.")
            return jsonify({"response": resp, "history": get_updated_history("assistant", resp)})
            
        # Route 2: Document Analysis
        if category == "Document Analysis":
            logger.info("Executing Document Analysis routing path")
            target_file_ids = []
            if selected_file_ids:
                target_file_ids = selected_file_ids
            elif intent == "compare":
                index_data = load_json_file(f"drive_index_{sanitized_email}.json", {})
                non_folders = [fid for fid in index_data if not index_data[fid].get("is_folder")]
                target_file_ids = non_folders[:2]
            elif selected_file_id:
                target_file_ids = [selected_file_id]
            else:
                resolved_file = resolve_file_by_name(user_email, message)
                if resolved_file:
                    target_file_ids = [resolved_file["id"]]
                    
            if not target_file_ids:
                resp = "⚠️ Please select a file or specify a filename to execute this analysis tool."
                return jsonify({"response": resp, "history": get_updated_history("assistant", resp)})
                
            file_records = get_file_records(user_email, service, target_file_ids)
            if not file_records:
                resp = "⚠️ Failed to extract content for the target files."
                return jsonify({"response": resp, "history": get_updated_history("assistant", resp)})
                
            tool_map = {
                "extract text": "extract_text",
                "summarize": "summarizer",
                "document summary": "summarizer",
                "analyze": "summarizer",
                "extract entities": "extract_entities",
                "entities": "extract_entities",
                "insights": "deep_insights",
                "compare": "compare_documents",
                "simplify": "simplify_document",
                "translate": "translate_document",
                "timeline": "extract_timeline",
                "sentiment": "sentiment_analysis",
                "citation": "citation_extraction",
                "takeaways": "important_points"
            }
            tool_name = tool_map.get(intent, "summarizer")
            
            kwargs = {}
            if tool_name == "translate_document":
                target_lang = "Spanish"
                for lang in ["spanish", "french", "german", "italian", "portuguese", "hindi", "telugu", "tamil", "chinese", "japanese", "russian"]:
                    if lang in message.lower():
                        target_lang = lang.title()
                        break
                kwargs["target_language"] = target_lang
                
            try:
                logger.info(f"Calling document_tools for tool: {tool_name}")
                resp = document_tools.run_tool(tool_name, file_records, **kwargs)
                
                if len(file_records) == 1:
                    fid = file_records[0]["id"]
                    cache_mapping = {
                        "summarizer": "summary",
                        "extract_entities": "entities",
                        "important_points": "important_points"
                    }
                    cache_key = cache_mapping.get(tool_name)
                    if cache_key:
                        val = file_records[0].get("generated_summary") if tool_name == "summarizer" else None
                        if not val:
                            val = resp
                        save_document_analysis(user_email, fid, cache_key, val)
                        
                names = [rec["name"] for rec in file_records]
                resolved_prefix = f"**Analyzing Files**: {', '.join(names)}\n\n"
                final_resp = resolved_prefix + resp
                return jsonify({"response": final_resp, "history": get_updated_history("assistant", final_resp)})
            except Exception as e_run:
                logger.exception("Document Analysis Failure")
                return jsonify({
                    "success": True,
                    "response": "AI service is temporarily unavailable or Try Again ",
                    "history": get_updated_history("assistant", "AI service is temporarily unavailable.")
                })
                
        # Route 3: Document Q&A
        if category == "Document Q&A":
            logger.info("Executing Document Q&A routing path")
            
            active_fid = selected_file_id
            if not active_fid and selected_file_ids:
                active_fid = selected_file_ids[0]
            if not active_fid:
                resolved_file = resolve_file_by_name(user_email, message)
                if resolved_file:
                    active_fid = resolved_file["id"]
                    
            if not active_fid:
                resp = "⚠️ Please select a file or specify a filename to ask questions about it."
                return jsonify({"response": resp, "history": get_updated_history("assistant", resp)})
                
            doc = get_document_content(user_email, service, active_fid)
            if not doc:
                resp = "⚠️ Failed to extract content for the target file."
                return jsonify({"response": resp, "history": get_updated_history("assistant", resp)})
                
            filename = doc.get("file_name", "Document")
            text = doc.get("full_text", "")
            word_count = len(text.split())
            summary = doc.get("summary", "")
            
            context_parts = [f"### [Active Document: {filename} (ID: {active_fid})]"]
            
            chunks = []
            if word_count < 3000:
                context_parts.append(f"Full Document text contents:\n\"\"\"\n{text}\n\"\"\"\n")
            else:
                if summary:
                    context_parts.append(f"Cached Summary Overview:\n{summary}\n")
                    
                try:
                    logger.info("Starting semantic search")
                    query_embs = None
                    if api_key:
                        query_embs = embeddings_db.get_embeddings_batch([message], api_key, provider=provider, is_query=True)
                    if query_embs:
                        matches = embeddings_db.semantic_search(
                            user_email, query_embs[0], limit=5, file_id=active_fid
                        )
                        chunks = [m[0]['text'] for m in matches]
                except Exception as e_vdb:
                    logger.exception("Vector Search Failure")
                    chunks = []
                    
                if not chunks:
                    doc_chunks = doc.get("chunks", [])
                    if not doc_chunks and text:
                        doc_chunks = chunk_text(text)
                    chunks = doc_chunks[:3]
                    context_parts.append("Retrieval Note: Loaded first fallback text chunks.")
                    
                context_parts.append("Relevant Document Excerpts:")
                for idx, chunk in enumerate(chunks):
                    context_parts.append(f"Excerpt {idx+1}:\n\"\"\"\n{chunk}\n\"\"\"\n")
                    
            context_str = "\n".join(context_parts)
            
            system_instruction = (
                "You are the Drive AI Assistant, a helpful assistant operating over the user's selected document.\n"
                "Answer the user's question in clear Markdown based strictly on the retrieved document contents and excerpts provided.\n"
                "Prioritize the active file contents and answer the question ONLY from them.\n"
                "Always cite the file name when referring to documents."
            )
            
            prompt = (
                f"{system_instruction}\n\n"
                f"### Retrieved Document Context (RAG):\n"
                f"{context_str}\n"
                f"### User Message:\n"
                f"\"{message}\"\n\n"
                "Please generate your response."
            )
            
            try:
                response_text = None
                if api_key:
                    logger.info("Calling LLM")
                    headers = {"Authorization": f"Bearer {nvidia_key}", "Content-Type": "application/json"}
                    payload = {
                        "model": "meta/llama-3.1-8b-instruct",
                        "messages": [
                            {"role": "system", "content": "You are the Drive AI Assistant."},
                            {"role": "user", "content": prompt}
                        ],
                        "temperature": 0.2,
                        "max_tokens": 1000
                    }
                    res = requests.post("https://integrate.api.nvidia.com/v1/chat/completions", headers=headers, json=payload, timeout=25)
                    if res.status_code == 200:
                        response_text = res.json()['choices'][0]['message']['content']
                    else:
                        raise Exception(f"Nvidia API returned code {res.status_code}: {res.text}")
                        
                if response_text:
                    cleaned_response = clean_assistant_response(response_text)
                    resolved_prefix = f"**Active File Context**: {filename}\n\n"
                    final_resp = resolved_prefix + cleaned_response
                    return jsonify({"response": final_resp, "history": get_updated_history("assistant", final_resp)})
                else:
                    if not api_key:
                        resp = "⚠️ No configured LLM API Key (Nvidia or Gemini) was found. Please set your API key by typing `set key [YOUR_API_KEY]`."
                        return jsonify({"response": resp, "history": get_updated_history("assistant", resp)})
                    raise Exception("Empty response text from LLM service")
            except Exception as e_qa:
                logger.exception("LLM Failure")
                return jsonify({
                    "success": True,
                    "response": "AI service is temporarily unavailable. or Try Again ",
                    "history": get_updated_history("assistant", "AI service is temporarily unavailable.")
                })
                
        # Route 4: General AI Chat
        if category == "General AI Chat":
            logger.info("Executing General AI Chat routing path")
            folder_contents = get_folder_contents_desc(user_email, current_folder_id)
            system_instruction = (
                "You are the Drive AI Assistant, a helpful conversational AI assistant.\n"
                "Answer the user's question directly, clearly, and helpfully in structured Markdown."
            )
            prompt = (
                f"{system_instruction}\n\n"
                f"### Current Folder Context:\n"
                f"Viewing Folder ID: `{current_folder_id}`\n"
                f"Contents of this folder:\n{folder_contents}\n\n"
                f"### User Message:\n"
                f"\"{message}\"\n\n"
                "Please generate your response."
            )
            
            try:
                response_text = None
                if api_key:
                    logger.info("Calling LLM")
                    headers = {"Authorization": f"Bearer {nvidia_key}", "Content-Type": "application/json"}
                    payload = {
                        "model": "meta/llama-3.1-8b-instruct",
                        "messages": [
                            {"role": "system", "content": "You are the Drive AI Assistant."},
                            {"role": "user", "content": prompt}
                        ],
                        "temperature": 0.3,
                        "max_tokens": 800
                    }
                    res = requests.post("https://integrate.api.nvidia.com/v1/chat/completions", headers=headers, json=payload, timeout=25)
                    if res.status_code == 200:
                        response_text = res.json()['choices'][0]['message']['content']
                    else:
                        raise Exception(f"Nvidia API returned code {res.status_code}: {res.text}")
                        
                if response_text:
                    cleaned_response = clean_assistant_response(response_text)
                    return jsonify({"response": cleaned_response, "history": get_updated_history("assistant", cleaned_response)})
                else:
                    if not api_key:
                        resp = "⚠️ No configured LLM API Key (Nvidia or Gemini) was found. Please set your API key by typing `set key [YOUR_API_KEY]`."
                    elif "hello" in msg_lower or "hi" in msg_lower:
                        resp = "🤖 Hello! I am the Drive AI Assistant. How can I help you manage or analyze your Google Drive today?"
                    elif "help" in msg_lower:
                        resp = "🤖 **How to use Drive AI Assistant**:\n- Ask general questions directly.\n- Analyze a document by selecting it and typing `Summarize this file`.\n- Manage folders and files by typing `Create folder Notes`."
                    else:
                        raise Exception("Empty response text from LLM service")
                    return jsonify({"response": resp, "history": get_updated_history("assistant", resp)})
            except Exception as e_gen:
                logger.exception("LLM Failure")
                return jsonify({
                    "success": True,
                    "response": "AI service is temporarily unavailable or Try Again ",
                    "history": get_updated_history("assistant", "AI service is temporarily unavailable.")
                })
                
        # Default chatbot help message fallback
        resp = (
            "🤖 **Drive AI Assistant**\n\n"
            "I can reason over your entire Google Drive! Ask me questions about your documents, compare files, or trigger file management commands:\n"
            "- `create folder Projects` - Create a folder.\n"
            "- `create doc Ideas` - Create a text document.\n"
            "- `delete report` - Resolve report and move it to trash.\n"
            "- `star notes` - Star notes file.\n"
            "- `rename document to project_plan` - Rename a file.\n\n"
            "**Semantic Search & Q&A**:\n"
            "- *\"Find all files discussing machine learning embeddings\"*\n"
            "- *\"Summarize what the marketing brief says about the budget\"*\n"
            "- *\"Are there any project milestones mentioned in planning.pptx?\"*\n\n"
            "🔑 *Tip: Paste your Nvidia API key as `set key [API_KEY]` to enable full LLM capabilities!*"
        )
        return jsonify({"response": resp, "history": get_updated_history("assistant", resp)})
    except Exception as e:
        logger.exception("AI Chat Fatal Error")
        return jsonify({
            "success": True,
            "response": "AI service is temporarily unavailable or Try Again"
        })
     


if __name__ == '__main__':
    # Run server on port 5000 to match the OAuth redirect URI in client_secrets.json.
    # Debug mode is OFF by default for safety; set FLASK_DEBUG=1 for local development.
    debug_mode = os.environ.get('FLASK_DEBUG', '0') == '1'
    app.run(debug=debug_mode, port=int(os.environ.get('PORT', 5000)))

